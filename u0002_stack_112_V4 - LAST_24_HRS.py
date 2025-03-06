#!/usr/bin/env python
# coding: utf-8
## this version is for the new 112 & dummy stack placed in spiarare system.
# In[4]:

# import the relevant libraries
import datetime
from datetime import datetime
from plotly.subplots import make_subplots
import plotly.subplots as sp
import plotly.graph_objects as go
import plotly.express as px
import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from datetime import date
from pptx.dml.color import RGBColor
import os
import sys
from scipy import stats

# Import the required tkinter Libraries
import tkinter as tk
from tkinter import *
from tkinter import ttk, filedialog

#from tkinter.filedialog import askopenfile
from tkinter import messagebox
from threading import Thread
import zipfile

## Image libraries
import plotly.io as pio
import tempfile
from PIL import Image
import warnings
warnings.filterwarnings('ignore')

# In[5]:

##### L1 ######
List_active_cells = [
    'C1-1', 'C1-2', 'C1-3', 'C1-4', 'C1-5', 'C1-6', 'C1-7', 'C1-8', 'C1-10', 'C1-11', 'C1-12', 'C1-13', 'C1-14', 'C1-15', 'C1-16', 'C2-1', 'C2-3', 'C2-4', 'C2-5', 'C2-6', 'C2-7', 'C2-8', 'C2-9', 'C2-10', 'C2-12', 'C2-13', 'C2-14', 'C2-15', 'C2-16', 'C3-1', 'C3-2', 'C3-3', 'C3-5', 'C3-6', 'C3-7', 'C3-8', 'C3-9', 'C3-10', 'C3-11', 'C3-12', 'C3-14', 'C3-15', 'C3-16', 'C4-1', 'C4-2', 'C4-3', 'C4-4', 'C4-5', 'C4-7', 'C4-8'
                    ,'C4-9', 'C4-10', 'C4-11', 'C4-12', 'C4-13', 'C4-14',  'C4-16'
                  ,'C5-1', 'C5-2', 'C5-3', 'C5-4', 'C5-5', 'C5-6', 'C5-7', 'C5-9', 'C5-10', 'C5-11', 'C5-12', 'C5-13', 'C5-14', 'C5-15', 'C5-16',  'C6-2', 'C6-3', 'C6-4', 'C6-5', 'C6-6', 'C6-7', 'C6-8', 'C6-9', 'C6-11', 'C6-12', 'C6-13', 'C6-14', 'C6-15', 'C6-16', 'C7-1', 'C7-2', 'C7-4', 'C7-5', 'C7-6', 'C7-7', 'C7-8', 'C7-9', 'C7-10', 'C7-11', 'C7-13', 'C7-14', 'C7-15', 'C7-16'       
                                      ,'C8-1', 'C8-2', 'C8-3', 'C8-4', 'C8-6', 'C8-7', 'C8-8'
                                    , 'C8-9', 'C8-10', 'C8-11', 'C8-12', 'C8-13', 
#last three cells               'C8-14', 'C8-15', 'C8-16',
#every 9th cell in the stack    'C1-9','C2-2','C2-11','C3-4','C3-13','C4-6','C4-15','C5-8','C6-1','C6-10','C7-3','C7-12','C8-5'
                  'C9-1', 'C9-2', 'C9-3', 'C9-4', 'C9-5', 'C9-6', 'C9-7', 'C9-8',  'C9-10', 'C9-11', 'C9-12', 'C9-13', 'C9-14', 'C9-15', 'C9-16', 'C10-1',  'C10-3', 'C10-4', 'C10-5', 'C10-6', 'C10-7', 'C10-8', 'C10-9', 'C10-10',  'C10-12', 'C10-13', 'C10-14', 'C10-15', 'C10-16', 'C11-1', 'C11-2', 'C11-3',  'C11-5', 'C11-6', 'C11-7', 'C11-8', 'C11-9', 'C11-10', 'C11-11', 'C11-12',  'C11-14', 'C11-15', 'C11-16', 'C12-1', 'C12-2', 'C12-3', 'C12-4', 'C12-5',  'C12-7', 'C12-8',
                                'C12-9', 'C12-10', 'C12-11', 'C12-12', 'C12-13', 'C12-14','C12-16',
                                'C13-1', 'C13-2', 'C13-3', 'C13-4', 'C13-5', 'C13-6', 'C13-7', 'C13-9', 'C13-10', 'C13-11', 'C13-12', 'C13-13', 'C13-14', 'C13-15', 'C13-16',  'C14-2', 'C14-3', 'C14-4', 'C14-5', 'C14-6', 'C14-7', 'C14-8', 'C14-9',  'C14-11', 'C14-12', 'C14-13', 'C14-14', 'C14-15', 'C14-16', 'C15-1', 'C15-2',  'C15-4', 'C15-5', 'C15-6', 'C15-7', 'C15-8', 'C15-9', 'C15-10', 'C15-11',  'C15-13', 'C15-14', 'C15-15', 'C15-16', 'C16-1', 'C16-2', 'C16-3', 'C16-4',  'C16-6', 'C16-7', 'C16-8',
                                'C16-9', 'C16-10', 'C16-11', 'C16-12', 'C16-13'
 #last three cells                # 'C16-14', 'C16-15', 'C16-16'
 #every 9th cell in the stack     #'C9-9', 'C10-2', 'C10-11', 'C11-4', 'C11-13', 'C12-6','C12-15', 'C13-8','C14-1','C14-10','C15-12','C15-12','C16-5'
                                              ]

Stack_voltage_list = List_active_cells.copy()

# In[6]:

#Stack_voltage_list = List_C1_1_C8_16.copy()
#Stack_voltage_list = List_C9_1_C16_16.copy()
# In[7]:

def create_Vcell_data(sys_data):
    # selecting only voltage cells
    selected_columns = list(sys_data.loc[:, 'C1-1':'C16-16'].columns) + ['Date','STACK CURRENT','CURRENT SP','DURATION_HR']

    # Reassigning Vcell_data to include selected columns from sys_data
    Vcell_data = sys_data[selected_columns]
     
    # Drop rows where all the specified columns are blank
    Vcell_data.dropna(subset=list(sys_data.loc[:, 'C1-1':'C16-16'].columns), how='all', inplace=True)  
    #Vcell_data.dropna(how='all', inplace=True)
    Vcell_data = Vcell_data.fillna(0) # for a row where only blank and 0s are present
    
    
    # Drop rows where all specified columns are 0
    cols_to_check = sys_data.loc[:, 'C1-1':'C16-16'].columns
    Vcell_data = Vcell_data[(Vcell_data[cols_to_check] != 0).any(axis=1)]
    
# ######## calculating std_dev (use this if stack position changed/added/removed in between data)
#     columns_to_check_list = []
#     for idx, row in Vcell_data.iterrows():
#         columns_to_check = [col for col in Vcell_data.columns if row[col] != 0]
#         columns_to_check_list.append(columns_to_check)
#     print(Vcell_data.head())    
#     std_dev_list = []
#     for columns_to_check in columns_to_check_list:
#         std_dev = sys_data[cols_to_check].std(axis=1)
#         std_dev_list.append(std_dev)   
        
#     sys_data['std_dev'] = std_dev_list
    

    # In[6]:  
    #convert all negative voltage values to 0
    Vcell_data[cols_to_check] = Vcell_data[cols_to_check].apply(lambda row: row.clip(lower=0), axis=1)
    # Add a new column 'Avg_V' to store the average of non-zero values in each row
    Vcell_data['Avg_V'] = Vcell_data[Stack_voltage_list].sum(axis=1) / Vcell_data[Stack_voltage_list].apply(lambda row: row[row != 0].count(), axis=1)
    Vcell_data['Max_V'] = Vcell_data[Stack_voltage_list].max(axis=1) 
    Vcell_data['std_dev'] = Vcell_data[Stack_voltage_list].std(axis=1)

    # Calculate Min_V excluding zero values 
    Vcell_data_nonzero = Vcell_data[cols_to_check].replace(0, np.nan)
    Vcell_data['Min_V'] = Vcell_data_nonzero.min(axis=1) 
   
    ## making cd, voltage charts., and ASR charts with dataframe created - Vcell_data
    Vcell_data['CURRENT DENSITY'] = Vcell_data['STACK CURRENT'] / 702
    
    #ASR in mohm-cm2
    Vcell_data_asr = Vcell_data[Vcell_data['STACK CURRENT'] != 0]
    Vcell_data_asr['ASR_Stack1'] = (Vcell_data['Avg_V']/Vcell_data['CURRENT DENSITY'])*1000


    Vcell_data_asr['ASR(max)_Stack1'] = (Vcell_data['Max_V']/Vcell_data['CURRENT DENSITY'])*1000
  
    return Vcell_data

    
# In[8]:
# plot current density
def plot_cd(sys_data, image_folder_name):
    #avg_cd_per_hour = Vcell_data.loc[:,['Date','CURRENT DENSITY']].resample('H', on='Date').mean().reset_index()
    Current_Density = go.Scattergl(
                           x = sys_data['Date'],
                           y = sys_data['CURRENT DENSITY'],
                           mode = 'markers',
                           marker=dict(color='rgb(214, 39, 40)',size=2),
                           #showlegend = True,
    #                        mode = 'lines',
                           name = 'Current Density (A/c'+'m\u00b2)'
                           #text = df.column (this is the info displayed on hover)
                           )

    data_plot = [Current_Density]
    
    # configure the layout
    layout = go.Layout(
                #to define axis, can either use go.layout.etc or dict(etc)
                xaxis = go.layout.XAxis(
    #                                     title = 'Date',
                                        nticks= 4,
                                        #tickformat = "%d/%m/%y"+'<br>'+"  %H:%M",
                                        tickformat = "%d-%b"+'<br>'+" %H:00",
                                        #hoverformat = "%d/%m/%y"+'<br>'+"  %H:%M:%S",
                                        hoverformat = "%d-%b %H:00",
                                        showgrid = False,
                                        #range=[100,steady_state['Duration(hr)'].max()]
                                        ),    
                yaxis= dict(title='Current Density'+'<br>'+'(A/c'+'m\u00b2)',
                                      #overlaying='y',
                                     showgrid = False,
                                     #range = [0,6],
                                     rangemode='tozero'),
                #legend= dict(xanchor="right",x=1.2,yanchor="top", y =1.2, bgcolor = 'rgba(0,0,0,0)'),
                template = 'simple_white',
                font = dict(size = 18)
                )


    # create the figure
    fig0 = go.Figure(data = data_plot, layout = layout)
    #fig0.update_xaxes(tickangle=90)
    #fig0.update_yaxes(title="Current Density'+'<br>'+'(A/c'+'m\u00b2)")
    fig0.update_layout(xaxis = dict(titlefont = dict(size = 18), tickfont = dict(size = 18), tickformat = "%d-%b"+'<br>'+" %H:00"),
           yaxis = dict(titlefont = dict(size = 18), tickfont = dict(size = 18)),
           yaxis2 = dict(titlefont = dict(size = 18), tickfont = dict(size = 18)) 
                     )
    #fig0.update_layout(yaxis_title="Current Density'+'<br>'+'(A/c'+'m\u00b2)")
    fig0.write_image('{}/Current_Density.png'.format(image_folder_name), engine="orca")
    fig0.write_html('{}/Current_Density.html'.format(image_folder_name), include_plotlyjs=('cdn'))
#     fig0.show()
    print('exported Current_Density plot')


# In[9]:
    
# defining function for avg stack voltage chart
def avg_min_max_stack_v(sys_data,image_folder_name):
    mean_trace = go.Scatter(
            x = sys_data.Date,
            y = sys_data['Avg_V_C1_1_C8_16'],
            mode = 'lines',
            name = 'Avg V',
            #line = dict(color = 'rgb(31, 119, 180)', shape='hv'),
            line=dict(color='rgb(31, 119, 180)', shape='hv', width=0.5),
            connectgaps=False
            )
    max_trace = go.Scatter(
        x=list(sys_data.Date), 
        y=list(sys_data['Max_V_C1_1_C8_16']), 
        line=dict(width=0, shape = 'hv', color = 'red'),
        mode='lines',
        name = 'max',
        showlegend=False,
        connectgaps=False
        )
    min_trace = go.Scatter(
        name='Min',
        x=list(sys_data.Date),
        y=list(sys_data['Min_V_C1_1_C8_16']),
        mode='lines',
        fillcolor='rgba(31,119,180,0.2)',
        line=dict(width = 0, shape = 'hv', color = 'green'),
        fill='tonexty',
        showlegend=False,
        connectgaps=False
    )
     
    std_dev_trace = go.Scatter(
        x = sys_data.Date,
        y = sys_data['std_dev'],
        mode = 'lines',
        name = 'Std Dev',
        #line = dict(color = 'rgb(255, 127, 14)', shape='hv'),
        line=dict(color='rgb(255, 127, 14)', shape='hv', width=0.5),
        connectgaps=False
        )
     
    dashed_line = {
            'type': 'line',
            'xref': 'paper',
            'yref': 'y',
            'x0': 0,
            'y0': 2,
            'x1': 1,
            'y1': 2,
            'line': {
                'color': 'black',
                'width': 2,
                'dash': 'dash',
            },
        }
    data = [mean_trace, max_trace, min_trace, std_dev_trace]
     
    layout = go.Layout(
                #title = 'plot_title',
                #to define axis, can either use go.layout.etc or dict(etc)
                xaxis = go.layout.XAxis(
          #                                   title = 'Date',
                                              nticks= 4,
                                              #tickformat = "%d/%m/%y"+'<br>'+"  %H:%M",
                                              tickformat = "%d-%b"+'<br>'+"%H:00",
                                              #hoverformat = "%d/%m/%y"+'<br>'+"  %H:%M:%S",
                                              hoverformat = "%d-%b %H:00",
                                              showgrid = False,
                                              #range=[100,steady_state['Duration(hr)'].max()]
                                              ),
                yaxis = dict(title='Voltage (V)', range=[1.4, 2]),
                #             showgrid = False),
                # legend= dict(x=0, y =1.1, bgcolor = 'rgba(0,0,0,0)'),
                # legend_title = None,\
                template = 'simple_white',
                shapes=[dashed_line],
                font = dict(size = 18)
                )
        
            
        
    fig = go.Figure(data = data, layout = layout)
    fig.show()    
    fig.write_image('{}/LHC_Stack_Voltage_&_Current_Density.png'.format(image_folder_name), engine="orca")
    fig.write_html('{}/LHC_Stack_Voltage_&_Current_Density.html'.format(image_folder_name), include_plotlyjs=('cdn'))

    print('exported LHC_Stack_Voltage chart')     
   
#%%
def avg_min_max_stack_v_old(sys_data, image_folder_name):

    #avg_per_hour = Vcell_data.loc[:,['Date','CURRENT DENSITY','Avg_V_C1_1-C8_16', 'Max_V_C1_1-C8_16', 'Min_V_C1_1-C8_16']].resample('H', on='Date').mean().reset_index()

    # define a new figure
    fig = sp.make_subplots(rows=4, cols=1, shared_xaxes=True, vertical_spacing=0.03)
    fig.update_layout(
                #title = 'LHC Cell Voltage & Current Density',
                #to define axis, can either use go.layout.etc or dict(etc)
                
                xaxis4 = dict(#title = 'Date',
                              nticks= 4,
                              tickformat = "%d-%b"+'<br>'+" %H:00",
                              #hoverformat = "%d/%m/%y"+'<br>'+"  %H:%M:%S",
                              hoverformat = "%d-%b %H:00",
                              showgrid = False),                                                                                
                yaxis = dict(title='Avg (V)',
    #                                font = 15,
                                    #range = [1.3, 2.3],
                                    showgrid = False                        
                                    #rangemode='tozero'
                                    ),
                yaxis2 = dict(title='Max (V)',
    #                                 font = 15,
                                    #range = [1.3, 2.3],
                                    showgrid = False
                                    #rangemode='tozero'
                                    ),
                yaxis3 = dict(title='Min (V)',
    #                                 font = 15,
                                    #range = [1.3, 2.3],
                                    showgrid = False
                                    #rangemode='tozero'
                                    ),
                yaxis4= dict(title='i'+' (A/c'+'m\u00b2)',
                                     #  range = [0,2],
                                     #overlaying='y',
                                     showgrid = False
                                     #rangemode='tozero'
                                     ),
                #legend= dict(xanchor="right",x=1.5,yanchor="top", y =2, bgcolor = 'rgba(0,0,0,0)'),               
                #legend= dict(xanchor="right",yanchor="top", y=1.38, bgcolor = 'rgba(0,0,0,0)'),
                legend= dict(bgcolor = 'rgba(0,0,0,0)'),
                template = 'simple_white',
                font = dict(size = 12)
                )

    fig.update_layout(xaxis = dict(titlefont = dict(size = 12), tickfont = dict(size = 12)),
           yaxis = dict(titlefont = dict(size = 12), tickfont = dict(size = 12)),
           yaxis2 = dict(titlefont = dict(size = 12), tickfont = dict(size = 12)),
           yaxis3 = dict(titlefont = dict(size = 12), tickfont = dict(size = 12)),
           yaxis4 = dict(titlefont = dict(size = 12), tickfont = dict(size = 12)) )
    color_dict = {'U1':fig.layout['template']['layout']['colorway'][0] ,
                  'U2':fig.layout['template']['layout']['colorway'][1] ,
                  'L2':fig.layout['template']['layout']['colorway'][2] ,
                  'L1':fig.layout['template']['layout']['colorway'][3] ,
                  'i' :fig.layout['template']['layout']['colorway'][4] } 
    
    # defining avg voltage traces
    Avg_StackL1 = go.Scattergl(
                           x = sys_data['Date'],
                           y = sys_data['Avg_V_C1_1_C8_16'],
                           mode = 'markers',
                           #mode = 'lines',
                           legendgroup=1,
                           name = 'L1 (V)',
                           marker=dict(color = color_dict['L1'])
                           #yaxis = 'y1'
                           #text = df.column (this is the info displayed on hover)
                           )
    
    #adding traces to the figure
    fig.add_trace(Avg_StackL1, row=1, col=1)
    
    # defining max voltage traces
    
    Max_StackL1 = go.Scattergl(
                           x = sys_data['Date'],
                           y = sys_data['Max_V_C1_1_C8_16'],
                           mode = 'markers',
                           #mode = 'lines',
                           name = 'L1 (V)',
                           legendgroup=1,
                           showlegend=False,
                           marker=dict(color = color_dict['L1'])
                           #yaxis = 'y1'
                           #text = df.column (this is the info displayed on hover)
                           )
    #adding traces to the figure
    fig.add_trace(Max_StackL1, row=2, col=1)
   
   # defining Min voltage traces 
    Min_StackL1 = go.Scattergl(
                           x = sys_data['Date'],
                           y = sys_data['Min_V_C1_1_C8_16'],
                           mode = 'markers',
                           #mode = 'lines',
                           name = 'L1 (V)',
                           legendgroup=1,
                           showlegend=False,
                           marker=dict(color = color_dict['L1'])
                           #yaxis = 'y1'
                           #text = df.column (this is the info displayed on hover)
                           )
    #adding traces to the figure
    fig.add_trace(Min_StackL1, row=3, col=1)
    
    # defining current density
    Current_Density = go.Scattergl(
                           x = sys_data['Date'],
                           y = sys_data['CURRENT DENSITY'],
                           mode = 'markers',
    #                        mode = 'lines',
                           name = 'i (A/c'+'m\u00b2)',
                           marker=dict(color = color_dict['i']),
                           #text = df.column (this is the info displayed on hover)
                           )
    #adding cd trace to the figure
    fig.add_trace(Current_Density, row=4, col=1)
    

    fig.write_image('{}/LHC_Stack_Voltage_&_Current_Density.png'.format(image_folder_name), engine="orca")
    fig.write_html('{}/LHC_Stack_Voltage_&_Current_Density.html'.format(image_folder_name), include_plotlyjs=('cdn'))
    # fig7.show()
    print('exported LHC_Stack_Voltage chart')
   
 
# In[12]:

# func for creating ASR chart
def ASR_chart(Vcell_data_asr, image_folder_name):
    # converting ASR values to avg_per_hour
    #avg_per_hour = Vcell_data.loc[:,['Date','CURRENT DENSITY','ASR_Stack1','ASR(max)_Stack1']].resample('H', on='Date').mean().reset_index()

    # define a new figure
    fig = sp.make_subplots(rows=3, cols=1, shared_xaxes=True, vertical_spacing=0.07)
    fig.update_layout(
                #title = 'Stack ASR Chart',
                xaxis3 = dict(
                              #title = 'Date',
                              nticks= 4,
                              tickformat = "%d-%b"+'<br>'+" %H:00",
                              #hoverformat = "%d/%m/%y"+'<br>'+"  %H:%M:%S",
                              hoverformat = "%d-%b %H:00",
                              showgrid = False
                              #range=[100,steady_state['Duration(hr)'].max()]
                             ),    
                yaxis = dict(title='ASR'+'<br>'+'(mΩ-c'+'m\u00b2)',
    #                                 font = 15,
                                    range = [0,10000],
                                    showgrid = False,
                                    rangemode='tozero'),
                yaxis2= dict(title='Max ASR' +'<br>'+'(mΩ-c'+'m\u00b2)',
                                    range = [0,10000],
    #                                 font = 15,
                                    showgrid = False,
                                    rangemode='tozero'),
                yaxis3= dict(title='i'+'<br>'+'(A/c'+'m\u00b2)',
                                     #overlaying='y',
                                     #side = 'right',
                                     showgrid = False,
                                     rangemode='tozero',
                                     range=[0,2]),
                #legend= dict(xanchor="right",x=1.5,yanchor="top", y=2, bgcolor = 'rgba(0,0,0,0)'),
                #legend= dict(xanchor="right",yanchor="top", y=1.38, bgcolor = 'rgba(0,0,0,0)'),
                legend= dict(bgcolor = 'rgba(0,0,0,0)'),
                template = 'simple_white',
                font = dict(size = 12)
                )

    fig.update_layout(xaxis = dict(titlefont = dict(size = 12), tickfont = dict(size = 12)),
           yaxis = dict(titlefont = dict(size = 12), tickfont = dict(size = 12)),
           yaxis2 = dict(titlefont = dict(size = 12), tickfont = dict(size = 12)),
           yaxis3 = dict(titlefont = dict(size = 12), tickfont = dict(size = 12)) )
    
    color_dict = {'U1':fig.layout['template']['layout']['colorway'][0] ,
                  'U2':fig.layout['template']['layout']['colorway'][1] ,
                  'L2':fig.layout['template']['layout']['colorway'][2] ,
                  'L1':fig.layout['template']['layout']['colorway'][3] ,
                  'i' :fig.layout['template']['layout']['colorway'][4] }    
    
    ASR_StackU1 = go.Scattergl(
                           x = Vcell_data_asr['Date'],
                           y = Vcell_data_asr['ASR_Stack1'],
                           mode = 'markers',
                           #mode = 'lines',
                           name = 'L1 (V)',
                           legendgroup=1,
                           marker=dict(color = color_dict['L1'])
                           #yaxis = 'y1'
                           #text = df.column (this is the info displayed on hover)
                           )
    # ASR_StackU2 = go.Scattergl(
    #                        x = avg_per_hour['date_&_hrs'],
    #                        y = avg_per_hour['ASR_Stack2'],
    #                        mode = 'markers',
    #                        #mode = 'lines',
    #                        name = 'U2 (V)',
    #                        legendgroup=2,
    #                        marker=dict(color = color_dict['U2'])
    #                        #yaxis = 'y1'
    #                        #text = df.column (this is the info displayed on hover
    #                        )
    # ASR_StackL2 = go.Scattergl(
    #                        x = avg_per_hour['date_&_hrs'],
    #                        y = avg_per_hour['ASR_Stack3'],
    #                        mode = 'markers',
    #                        #mode = 'lines',
    #                        name = 'L2 (V)',
    #                        legendgroup=3,
    #                        marker=dict(color = color_dict['L2'])
    #                        #yaxis = 'y1'
    #                        #text = df.column (this is the info displayed on hover)
    #                        )
    # ASR_StackL1 = go.Scattergl(
    #                        x = avg_per_hour['date_&_hrs'],
    #                        y = avg_per_hour['ASR_Stack4'],
    #                        mode = 'markers',
    #                        #mode = 'lines',
    #                        name = 'L1 (V)',
    #                        legendgroup=4,
    #                        marker=dict(color = color_dict['L1'])
    #                        #yaxis = 'y1'
    #                        #text = df.column (this is the info displayed on hover)
    #                        )
    #adding ASR traces to the figure
    fig.add_trace(ASR_StackU1, row=1, col=1)
    # fig.add_trace(ASR_StackU2, row=1, col=1)
    # fig.add_trace(ASR_StackL2, row=1, col=1)
    # fig.add_trace(ASR_StackL1, row=1, col=1)
    
    # defining Max Asr
    max_ASR_StackU1 = go.Scattergl(
                           x = Vcell_data_asr['Date'],
                           y = Vcell_data_asr['ASR(max)_Stack1'],
                           mode = 'markers',
                           #mode = 'lines',
                           name = 'L1 (V)',
                           legendgroup=1,
                           marker=dict(color = color_dict['L1']),
                           showlegend=False
                           #yaxis = 'y1'
                           #text = df.column (this is the info displayed on hover)
                           )
    # max_ASR_StackU2 = go.Scattergl(
    #                        x = avg_per_hour['date_&_hrs'],
    #                        y = avg_per_hour['ASR(max)_Stack2'],
    #                        mode = 'markers',
    #                        #mode = 'lines',
    #                        name = 'U2 (V)',
    #                        legendgroup=2,
    #                        showlegend=False,
    #                        marker=dict(color = color_dict['U2'])
    #                        #yaxis = 'y1'
    #                        #text = df.column (this is the info displayed on hover)
    #                        )
    # max_ASR_StackL2 = go.Scattergl(
    #                        x = avg_per_hour['date_&_hrs'],
    #                        y = avg_per_hour['ASR(max)_Stack3'],
    #                        mode = 'markers',
    #                        #mode = 'lines',
    #                        name = 'L2 (V)',
    #                        legendgroup=3,
    #                        showlegend=False,
    #                        marker=dict(color = color_dict['L2'])
    #                        #yaxis = 'y1'
    #                        #text = df.column (this is the info displayed on hover)
    #                        )
    # max_ASR_StackL1 = go.Scattergl(
    #                        x = avg_per_hour['date_&_hrs'],
    #                        y = avg_per_hour['ASR(max)_Stack4'],
    #                        mode = 'markers',
    #                        #mode = 'lines',
    #                        name = 'L1 (V)',
    #                        legendgroup=4,
    #                        showlegend=False,
    #                        marker=dict(color = color_dict['L1'])
    #                        #yaxis = 'y1'
    #                        #text = df.column (this is the info displayed on hover)
    #                        )
    #adding Max ASR traces to the figure
    fig.add_trace(max_ASR_StackU1, row=2, col=1)
    # fig.add_trace(max_ASR_StackU2, row=2, col=1)
    # fig.add_trace(max_ASR_StackL2, row=2, col=1)
    # fig.add_trace(max_ASR_StackL1, row=2, col=1)
    
    # defining current density
    Current_Density = go.Scattergl(
                           x = Vcell_data_asr['Date'],
                           y = Vcell_data_asr['CURRENT DENSITY'],
                           mode = 'markers',
    #                        mode = 'lines',
                           name = 'i (A/c'+'m\u00b2)',
                           marker=dict(color = color_dict['i'])
                           #text = df.column (this is the info displayed on hover)
                           )
    #adding cd trace to the figure
    fig.add_trace(Current_Density, row=3, col=1)

    fig.write_image('{}/Stack_ASR_Chart.png'.format(image_folder_name), engine="orca")
    fig.write_html('{}/Stack_ASR_Chart.html'.format(image_folder_name), include_plotlyjs=('cdn'))
#     fig.show()
    print('exported Stack_ASR_Chart')                

# In[14]:

# func for setting upper and lower limits of sys_data
def limit_setting(sys_data):
    sys_data['CURRENT DENSITY'] = sys_data['STACK CURRENT'] / 702
    sys_data['PMP_ON_OFF'] = np.where(sys_data['PMP 102'] == 'ON',1,0)

    Col_List_0to100 = ['PRT 702','PRT 703','PRT 704','PRT 705','PRT 706','KTC 901','HYS 101','HYS 401' ,'HYS 501' ,'HYS 102' ,'OXS 101' ,'LVL 101', 'TTC 102' ,'PRT 401' ,'PRT 402' ,'KTC 401','TTC 101' ,'TTC 301' ,'PRT 101']
    condition = (sys_data[Col_List_0to100] >= 0) & (sys_data[Col_List_0to100] <= 100)
    sys_data = sys_data[condition.all(axis=1)]

    sys_data = sys_data[sys_data['PRT 102'] >= 0]
    # sys_data = sys_data[sys_data['DPT 401'] >= -70]

    Col_List_0to200 = ['COS 701','COS 702','COS 101']
    condition = (sys_data[Col_List_0to200] >= 0) & (sys_data[Col_List_0to200] <= 200)
    sys_data = sys_data[condition.all(axis=1)]
    
    list_cols = ['TTC 101','TTC 102', 'TTC 301','KTC 401','LVL 101']
    condition = (sys_data[list_cols] <= 100)
    sys_data = sys_data[condition.all(axis=1)]
    print('done with setting upper and lower limit of sys_data')
    return sys_data
# In[15]:
#func to create pressure and conductivity chart
def pressure_and_conductivity(sys_data, image_folder_name):
    avg_per_hour = sys_data.loc[:,['Date','PRT 102', 'PRT 401', 'COS 101']].resample('H', on='Date').mean().reset_index()

        # define a new figure
    Pressure_102 = go.Scattergl(
                           x = sys_data['Date'],
                           y = sys_data['PRT 102'],
                           mode = 'markers',
                           #mode = 'lines',
                           name = 'PRT 102',
                           yaxis = 'y1',
                           marker=dict(size=3.5)
                           #text = df.column (this is the info displayed on hover)
                           )
    Pressure_401 = go.Scattergl(
                           x = sys_data['Date'],
                           y = sys_data['PRT 401'],
                           mode = 'markers',
                           #mode = lines,
                           name = 'PRT 401',
                           yaxis = 'y1',
                           marker=dict(size=3.5)
                           #text = df.column (this is the info displayed on hover)
                           )
    Conductivity_101 = go.Scattergl(
                            x = sys_data['Date'],
                            y = sys_data['COS 101'],
                            mode = 'markers',
                            #mode = 'lines',
                            name = 'COS 101',
                            yaxis = 'y2',
                            marker=dict(size=3.5)
                            #text = df.column (this is the info displayed on hover)
                            )
    # update the data list
    data = [Pressure_102, Pressure_401, Conductivity_101]
    
    
    # configure the layout
    layout = go.Layout(
                title = 'LHC Pressure & Conductivity chart',
                #to define axis, can either use go.layout.etc or dict(etc)
                xaxis = go.layout.XAxis(
    #                                     title = 'Date',
                                        nticks= 5,
                                        tickformat = "%d-%b"+'<br>'+" %H:00",
                                        #hoverformat = "%d/%m/%y"+'<br>'+"  %H:%M:%S",
                                        hoverformat = "%d-%b %H:00",
                                        showgrid = False,
                                        #tickfont=dict(size=16, color='black')  # Adjust font properties as needed
                                        #tickangle=90
                                        #range=[100,steady_state['Duration(hr)'].max()]
                                        ),
                yaxis = dict(title='Pressure (bar)',
    #                                 font = 15,
                                    showgrid = False,
                                    rangemode='tozero'),
                yaxis2= dict(title='COS 101 (µS/cm)',
    #                                  overlaying='y',
                                     side = 'right',
                                     showgrid = False,
                                     rangemode='tozero',
    #                                  range=[0,12.5]
                            ),
              #legend= dict(xanchor="right",x=1.3, yanchor="top", y =1.5, bgcolor = 'rgba(0,0,0,0)'),
              legend= dict(xanchor="left",x=0, yanchor="top", y =1.1, bgcolor = 'rgba(0,0,0,0)', orientation='h'),
              template = 'simple_white',
              font = dict(size = 18),
              width=850
                  )
    # create the figure
    fig1 = go.Figure(data = data, layout = layout)

    
    fig1.update_layout(xaxis = dict(titlefont = dict(size = 18), tickfont = dict(size = 16)),
    yaxis = dict(titlefont = dict(size = 18), tickfont = dict(size = 18)) )
    
    
    fig1.write_image('{}/LHC_Pressure_&_Conductivity_Chart.png'.format(image_folder_name), engine="orca")
    fig1.write_html('{}/LHC_Pressure_&_Conductivity_Chart.html'.format(image_folder_name), include_plotlyjs=('cdn'))
    # fig1.show()
    print('exported LHC_Pressure_&_Conductivity Chart')


# In[16]:

# func for creating TTC,KTC & LVL sensors chart
def ttc_ktc_lvl_sensors(sys_data, image_folder_name):
    avg_per_hour = sys_data.loc[:,['Date','TTC 101', 'TTC 102', 'TTC 301', 'KTC 401', 'LVL 101']].resample('H', on='Date').mean().reset_index()
    
    TTC101 = go.Scattergl(
                           x = sys_data['Date'],
                           y = sys_data['TTC 101'],
                           mode = 'markers',
                           #marker=dict(size=3),
                           #mode = 'lines',
                           name = 'TTC 101',
                           yaxis = 'y1',
                           marker=dict(size=3.5)
                           #text = df.column (this is the info displayed on hover)
                           )
    TTC102 = go.Scattergl(
                           x = sys_data['Date'],
                           y = sys_data['TTC 102'],
                           mode = 'markers',
                           #marker=dict(size=3),
                           #mode = lines,
                           name = 'TTC 102',
                           yaxis = 'y1',
                           marker=dict(size=3.5)
                           #text = df.column (this is the info displayed on hover)
                           )
    TTC301  = go.Scattergl(
                            x = sys_data['Date'],
                            y = sys_data['TTC 301'],
                            mode = 'markers',
                            #marker=dict(size=3),
                            #mode = 'lines',
                            name = 'TTC 301',
                            yaxis = 'y1',
                            marker=dict(size=3.5)
                            #text = df.column (this is the info displayed on hover)
                            )
    KTC401  = go.Scattergl(
                            x = sys_data['Date'],
                            y = sys_data['KTC 401'],
                            mode = 'markers',
                            #marker=dict(size=3),
                            #mode = 'lines',
                            name = 'KTC 401',
                            yaxis = 'y1',
                            marker=dict(size=3.5)
                            #text = df.column (this is the info displayed on hover)
                            )
    LVL101  = go.Scattergl(
                            x = sys_data['Date'],
                            y = sys_data['LVL 101'],
                            mode = 'markers',
                            #marker=dict(size=3),
                            #mode = 'lines',
                            name = 'LVL 101',
                            yaxis = 'y2',
                            marker=dict(size=3.5)
                            #text = df.column (this is the info displayed on hover)
                            ) 



    # update the data list
    data = [TTC101, TTC102, TTC301, KTC401, LVL101]

    # configure the layout
    layout = go.Layout(
                title = 'LHC Temp & Level sensor chart',
                #to define axis, can either use go.layout.etc or dict(etc)
                xaxis = go.layout.XAxis(
    #                                     title = 'Date',
                                        nticks= 5,
                                        tickformat = "%d-%b"+'<br>'+" %H:00",
                                        #hoverformat = "%d/%m/%y"+'<br>'+"  %H:%M:%S",
                                        hoverformat = "%d-%b %H:00",
                                        showgrid = False,
                                        #range=[100,steady_state['Duration(hr)'].max()]
                                        ),

                yaxis = dict(title='Temp(°C),',
                                    showgrid = False,
                                    range = [0,100],
                                    rangemode='tozero'),
                yaxis2= dict(title='Water level (%)',
                                     overlaying='y',
                                     side = 'right',
                                     showgrid = False,
                                     rangemode='tozero',
                                     range=[0,100]),
                legend= dict(xanchor="left",x=0,yanchor="top", y =1.1, bgcolor = 'rgba(0,0,0,0)', orientation='h'),
                template = 'simple_white',
                font = dict(size = 18),
                width=850
                )

    # create the figure
    fig2 = go.Figure(data = data, layout = layout)

    fig2.update_layout(xaxis = dict(titlefont = dict(size = 18), tickfont = dict(size = 18)),
           yaxis = dict(titlefont = dict(size = 18), tickfont = dict(size = 18)) )
    fig2.update_layout()
    
    fig2.write_image('{}/LHC_TEMP_&_Level_Sensor_Chart.png'.format(image_folder_name), engine="orca")
    fig2.write_html('{}/LHC_TEMP_&_Level_Sensor_Chart.html'.format(image_folder_name), include_plotlyjs=('cdn'))
    # fig2.show()
    print('exported LHC_TEMP_&_Level_Sensor Chart')

  
# In[17]:


# func for HYS 102 VS LVL 101 Chart
def HYS102_vs_LVL101(sys_data, image_folder_name):
    # temp_data = temp_data[temp_data['CURRENT SP'] >= 0]
    sys_data = sys_data[(sys_data['MP STATE']=='STACK PWR') & (sys_data['MS STATE']=='ARMED')]
    avg_per_hour = sys_data.loc[:,['Date','HYS 102', 'PRT 401','CURRENT DENSITY','LVL 101']].resample('H', on='Date').mean().reset_index()

    HYS102 = go.Scattergl(
                           x = sys_data['Date'],
                           y = sys_data['HYS 102'],
                           mode = 'markers',
                           #marker=dict(size=3),
                           #mode = 'lines',
                           name = 'HYS 102 (% LEL)',
                           yaxis = 'y1',
                           marker=dict(size=3.5,color='#ff7f0e')
                           #text = df.column (this is the info displayed on hover)
                           )

    LVL101  = go.Scattergl(
                            x = sys_data['Date'],
                            y = sys_data['LVL 101'],
                            mode = 'markers',
                            #marker=dict(size=3),
                            #mode = 'lines',
                            name = 'LVL 101 (%)',
                            yaxis = 'y1',
                            marker=dict(size=3.5, color='#1f77b4')
                            #text = df.column (this is the info displayed on hover)
                            )
    PRT401 = go.Scattergl(
                           x = sys_data['Date'],
                           y = sys_data['PRT 401'],
                           mode = 'markers',
                           #mode = lines,
                           name = 'PRT 401 (bar)',
                           yaxis = 'y1',
                           marker=dict(size=3.5, color='#d62728')
                           #text = df.column (this is the info displayed on hover)
                           )
    Current_Density = go.Scattergl(
                           x = sys_data['Date'],
                           y = sys_data['CURRENT DENSITY'],
                           mode = 'markers',
    #                        marker=dict(color='#e377c2'),
    #                        mode = 'lines',
                           name = 'Current Density (A/c'+'m\u00b2)',
                           yaxis = 'y2',
                           marker=dict(size=3.5, color='#2ca02c')
                           #text = df.column (this is the info displayed on hover)
                           )

    # update the data list
    data = [Current_Density, PRT401,LVL101,HYS102]

    # configure the layout
    layout = go.Layout(
                title = 'LHC HYS 102 Vs. LVL 101',
                #to define axis, can either use go.layout.etc or dict(etc)
                xaxis = go.layout.XAxis(
    #                                     title = 'Date',
                                        nticks= 4,
                                        tickformat = "%d-%b"+'<br>'+" %H:00",
                                        #hoverformat = "%d/%m/%y"+'<br>'+"  %H:%M:%S",
                                        hoverformat = "%d-%b %H:00",
                                        showgrid = False
                                        #range=[100,steady_state['Duration(hr)'].max()]
                                        ),
                yaxis = dict(title=u'H\u2082 (% LEL) leak &'+'<br>'+' Water level (%)',
                                    showgrid = False,
                                    rangemode='tozero'),
                yaxis2= dict(title='Current Density'+'<br>'+'(A/c'+'m\u00b2)',
                                     overlaying='y',
                                     side = 'right',
                                     showgrid = False,
                                     rangemode='tozero'),
                legend= dict(xanchor="right",x=1.3,yanchor="top", y =1.6, bgcolor = 'rgba(0,0,0,0)'),
                template = 'simple_white',
                font = dict(size = 18)
                )

    # create the figure
    fig3 = go.Figure(data = data, layout = layout)

    fig3.update_layout(xaxis = dict(titlefont = dict(size = 18), tickfont = dict(size = 18)),
           yaxis = dict(titlefont = dict(size = 18), tickfont = dict(size = 18)) )

    fig3.write_image('{}/LHC_HYS_102_VS_LVL_101.png'.format(image_folder_name), engine="orca")
    fig3.write_html('{}/LHC_HYS_102_VS_LVL_101.html'.format(image_folder_name), include_plotlyjs=('cdn'))
    # fig3.show()
    print('exported LHC_HYS_102_VS_LVL_101 chart')


# In[18]:
# func to create hydrogen & water sensor charts
def hyd_and_water_sensor(sys_data, image_folder_name):
    #sys_data = sys_data[(sys_data['MP STATE']=='STACK PWR') & (sys_data['MS STATE']=='ARMED')]
    #avg_per_hour = sys_data.loc[:,['Date','HYS 101', 'HYS 102', 'HYS 401','HYS 501','OXS 101','CURRENT DENSITY']].resample('H', on='Date').mean().reset_index()
    
    HYS101 = go.Scattergl(
                           x = sys_data['Date'],
                           y = sys_data['HYS 101'],
                           mode = 'markers',
                           #marker=dict(size=3),
                           #mode = 'lines',
                           name = 'HYS 101',
                           yaxis = 'y1',
                           marker=dict(size=3.5)
                           #text = df.column (this is the info displayed on hover)
                           )
    HYS102 = go.Scattergl(
                           x = sys_data['Date'],
                           y = sys_data['HYS 102'],
                           mode = 'markers',
                           #marker=dict(size=3),
                           #mode = 'lines',
                           name = 'HYS 102',
                           yaxis = 'y1',
                           marker=dict(size=3.5)
                           #text = df.column (this is the info displayed on hover)
                           )
    HYS401 = go.Scattergl(
                           x = sys_data['Date'],
                           y = sys_data['HYS 401'],
                           mode = 'markers',
                           #marker=dict(size=3),
                           #mode = lines,
                           name = 'HYS 401',
                           yaxis = 'y1',
                           marker=dict(size=3.5)
                           #text = df.column (this is the info displayed on hover)
                           )
    HYS501  = go.Scattergl(
                            x = sys_data['Date'],
                            y = sys_data['HYS 501'],
                            mode = 'markers',
                            #marker=dict(size=3),
                            #mode = 'lines',
                            name = 'HYS 501',
                            yaxis = 'y1',
                            marker=dict(size=3.5)
                            #text = df.column (this is the info displayed on hover)
                            )
    OXS101  = go.Scattergl(
                            x = sys_data['Date'],
                            y = sys_data['OXS 101'],
                            mode = 'markers',
                            #marker=dict(size=3),
                            #mode = 'lines',
                            name = 'OXS 101',
                            yaxis = 'y1',
                            marker=dict(size=3.5)
                            #text = df.column (this is the info displayed on hover)
                            )
    Current_Density = go.Scattergl(
                           x = sys_data['Date'],
                           y = sys_data['CURRENT DENSITY'],
                           mode = 'markers',
    #                        mode = 'lines',
                           name = 'Cd',
                           yaxis = 'y2',
                           marker=dict(size=3.5)
                           #text = df.column (this is the info displayed on hover)
                           )

    # update the data list
    data = [HYS101, HYS102, HYS401, HYS501, OXS101, Current_Density]

    # configure the layout
    layout = go.Layout(
                title = 'LHC '+u'H\u2082 & O\u2082'+' leak',
                #to define axis, can either use go.layout.etc or dict(etc)
                xaxis = go.layout.XAxis(
    #                                     title = 'Date',
                                        nticks= 5,
                                        tickformat ="%d-%b"+'<br>'+" %H:00",
                                        #hoverformat = "%d/%m/%y"+'<br>'+"  %H:%M:%S",
                                        hoverformat = "%d-%b %H:00",
                                        showgrid = False,
                                        #range=[100,steady_state['Duration(hr)'].max()]
                                        ),
                yaxis = dict(title=u'H\u2082 (% LEL) & O\u2082 (%) leak',
                                    showgrid = False,
                                    rangemode='tozero'),
                yaxis2= dict(title='Cd (A/c'+'m\u00b2)',
                                     overlaying='y',
                                     side = 'right',
                                     showgrid = False,
                                     range = [0,6],
                                     rangemode='tozero'),
                legend= dict(xanchor="left",x=0,yanchor="top", y =1.1, bgcolor = 'rgba(0,0,0,0)',orientation='h'),
                template = 'simple_white',
                font = dict(size = 18),
                width=850
                )

    # create the figure
    fig3 = go.Figure(data = data, layout = layout)

    fig3.update_layout(xaxis = dict(titlefont = dict(size = 18), tickfont = dict(size = 18)),
           yaxis = dict(titlefont = dict(size = 18), tickfont = dict(size = 18)) )

    fig3.write_image('{}/LHC_H2_&_O2_leak.png'.format(image_folder_name), engine="orca")
    fig3.write_html('{}/LHC_H2_&_O2_leak.html'.format(image_folder_name), include_plotlyjs=('cdn'))
    # fig3.show()
    print('exported hyd_and_water_sensor chart')
    
# In[19]:

# func to create LWC - PRT 704 and PMP state chart 
def LWC_Pressure_and_PMP_State(sys_data, image_folder_name):
    # define a new figure
    fig = sp.make_subplots(rows=2, cols=1, shared_xaxes=True, vertical_spacing=0.07)
    fig.update_layout(
                title = 'LWC Pressure & PMP state',
                xaxis2 = dict(
                              #title = 'Date',
                              nticks= 4,
                              tickformat = "%d-%b"+'<br>'+" %H:00",
                              #hoverformat = "%d/%m/%y"+'<br>'+"  %H:%M:%S",
                              hoverformat = "%d-%b %H:00",
                              showgrid = False
                              #range=[100,steady_state['Duration(hr)'].max()]
                             ),    
                yaxis = dict(title = 'PRT 704'+'<br>'+'(bar)',
                                    #range = [0,10000],
                                    showgrid = False,
                                    rangemode='tozero'),
                yaxis2 = dict(title = 'PRT 705'+'<br>'+'(bar)',
                                    #range = [0,10000],
                                    showgrid = False,
                                    rangemode='tozero'),
                yaxis3 = dict(title='PMP 102',
                                     #overlaying='y',
                                     side = 'right',
                                     tickvals=[0,1],
                                     showgrid = False,
                                     rangemode='tozero'
                                     ),
                yaxis4 = dict(title= 'PMP 102',
                                     #overlaying='y2',
                                     side = 'right',
                                     tickvals=[0,1],
                                     showgrid=False,
                                     rangemode='tozero'
                                     ),
                #legend = dict(xanchor="right",x=1.5,yanchor="top", y=2, bgcolor = 'rgba(0,0,0,0)'),
                legend= dict(xanchor="right",yanchor="top", y=1.38, bgcolor = 'rgba(0,0,0,0)'),
                #legend= dict(bgcolor = 'rgba(0,0,0,0)'),
                template = 'simple_white',
                font = dict(size = 12)
                )

    fig.update_layout(xaxis = dict(titlefont = dict(size = 12), tickfont = dict(size = 12)),
           yaxis = dict(titlefont = dict(size = 12), tickfont = dict(size = 12)),
           yaxis2 = dict(titlefont = dict(size = 12), tickfont = dict(size = 12)),
           yaxis3 = dict(titlefont = dict(size = 12), tickfont = dict(size = 12)),
           yaxis4 = dict(titlefont = dict(size = 12), tickfont = dict(size = 12)) )
    color_dict = {'PRT 704':fig.layout['template']['layout']['colorway'][0] ,
                  'PRT 705':fig.layout['template']['layout']['colorway'][4] ,
                  'PMP 102':fig.layout['template']['layout']['colorway'][1] } 
    
    fig.update_layout(yaxis3 = dict(side = 'right'),
                      yaxis4 = dict(side = 'right'))
    # defining LWC pressure
    PRT704 = go.Scattergl(
                           x = sys_data['Date'],
                           y = sys_data['PRT 704'],
                           mode = 'markers',
                           marker=dict(size=4, color = color_dict['PRT 704']),
    #                        mode = 'lines',
                           name = 'PRT 704 (bar)',
                           #yaxis = 'y1'
                           #text = df.column (this is the info displayed on hover)
                           )
    PRT705 = go.Scattergl(
                           x = sys_data['Date'],
                           y = sys_data['PRT 705'],
                           mode = 'markers',
                           marker=dict(size=4, color = color_dict['PRT 705']),
    #                        mode = 'lines',
                           name = 'PRT 705 (bar)',
                           #yaxis = 'y1'
                           #text = df.column (this is the info displayed on hover)
                           )
    # PRT706  = go.Scattergl(
    #                         x = sys_data['Date'],
    #                         y = sys_data['PRT 706'],
    #                         mode = 'markers',
    #                         #marker=dict(size=3),
    #                         #mode = 'lines',
    #                         name = 'PRT 706 (bar)',
    #                         yaxis = 'y1',
    #                         #text = df.column (this is the info displayed on hover)
    #                         )
    PMP102_a  = go.Scattergl(
                            x = sys_data['Date'],
                            y = sys_data['PMP_ON_OFF'],
                            mode = 'markers',
                            marker=dict(size=4, color = color_dict['PMP 102']),
                            legendgroup=1,
                            #mode = 'lines',
                            name = 'PMP State',
                            #showlegend=False,
                            yaxis = 'y3',
                            #text = df.column (this is the info displayed on hover)
                            )
    PMP102_b  = go.Scattergl(
                            x = sys_data['Date'],
                            y = sys_data['PMP_ON_OFF'],
                            mode = 'markers',
                            marker=dict(size=4, color = color_dict['PMP 102']),
                            legendgroup=1,
                            #mode = 'lines',
                            name = 'PMP State',
                            showlegend=False,
                            yaxis = 'y4',
                            #text = df.column (this is the info displayed on hover)
                            )
    
    #adding Max ASR traces to the figure
    fig.add_trace(PRT704, row=1, col=1)
    fig.add_trace(PRT705, row=2, col=1)
    fig.add_trace(PMP102_a, row=1, col=1)
    fig.add_trace(PMP102_b, row=2, col=1)
 
    fig.write_image('{}/LWC_Pressure_and_PMP_State.png'.format(image_folder_name), engine="orca")
    fig.write_html('{}/LWC_Pressure_and_PMP_State.html'.format(image_folder_name), include_plotlyjs=('cdn'))

    print('exported LWC_Pressure_and_PMP_State chart')


# In[21]:

# func to create LWC COS701 , COS702 Conductivity chart
def LWC_conductivity(sys_data, image_folder_name):
    COS701 = go.Scattergl(
                           x = sys_data['Date'],
                           y = sys_data['COS 701'],
                           mode = 'markers',
                           marker=dict(size=4),
                           #mode = 'lines',
                           name = 'COS 701 (µS/cm)',
                           yaxis = 'y1'
                           #text = df.column (this is the info displayed on hover)
                           )
    COS702 = go.Scattergl(
                           x = sys_data['Date'],
                           y = sys_data['COS 702'],
                           mode = 'markers',
                           marker=dict(size=4),
                           #mode = lines,
                           name = 'COS 702 (µS/cm)',
                           yaxis = 'y1'
                           #text = df.column (this is the info displayed on hover)
                           )
    # update the data list
    data = [COS701, COS702]

    # configure the layout
    layout = go.Layout(
                title = 'LWC Conductivity',
                #to define axis, can either use go.layout.etc or dict(etc)
                xaxis = go.layout.XAxis(
    #                                     title = 'Date',
                                        nticks= 4,
                                        tickformat = "%d-%b"+'<br>'+" %H:00",
                                        #hoverformat = "%d/%m/%y"+'<br>'+"  %H:%M:%S",
                                        hoverformat = "%d-%b %H:00",
                                        showgrid = False,
                                        #range=[100,steady_state['Duration(hr)'].max()]
                                        ),    
                yaxis = dict(title='Conductivity (µS/cm)',
                                    showgrid = False,
                                    rangemode='tozero'),
                yaxis2= dict(title='Current Density'+'<br>' +'(A/c'+'m\u00b2)',
                                     overlaying='y',
                                     side = 'right',
                                     range = [0,2],
                                     showgrid = False,
                                     rangemode='tozero'),
    #                                  range=[0,600]),
                legend = dict(xanchor="right",x=1.4,yanchor="top", y =1.3, bgcolor = 'rgba(0,0,0,0)'),
                template = 'simple_white',
                font = dict(size = 18)
                )

    # create the figure
    fig11 = go.Figure(data = data, layout = layout)

    fig11.update_layout(xaxis = dict(titlefont = dict(size = 18), tickfont = dict(size = 18)),
                      yaxis = dict(titlefont = dict(size = 18), tickfont = dict(size = 18)),
    #                   yaxis2 = dict(titlefont = dict(size = 15), tickfont = dict(size = 15)) 
                     )


    fig11.write_image('{}/LWC_Conductivity.png'.format(image_folder_name), engine="orca")
    fig11.write_html('{}/LWC_Conductivity.html'.format(image_folder_name), include_plotlyjs=('cdn'))
    # fig11.show()
    print('exported LWC_conductivity chart')


# In[22]:


# func to define t_start, t_end and N for polcurves
def t_start_t_end_N(Vcell_data): 
    # Drop rows where all the specified columns are blank
    # Vcell_data.dropna(subset=Stack_voltage_list, how='all', inplace=True)

    # make a copy of the relevant df
    temp = Vcell_data.loc[:,['Date','CURRENT DENSITY']].copy()
    # identify pol curve end
    # this is defined here as where the difference in rolling mean is < -.05 and current density is > than 0.5
    end_condition1 = (temp['CURRENT DENSITY'].rolling(5, min_periods=1, center=True).mean() - temp['CURRENT DENSITY'].shift().rolling(5, min_periods=1, center=True).mean()) < -0.01
    # print(end_condition1.value_counts())
    end_condition2 = (temp['CURRENT DENSITY']) > 0.1
    # print(end_condition2.value_counts())
    end_condition = end_condition1.shift(-1, fill_value=False) & end_condition2
    endpoints = temp.loc[end_condition, :].copy()
    endpoints['type'] = 1
    # print(endpoints)
    # idnetify pol curve starts
    # this is defined here as where the current density increases by > 0.001 and is less than 0.02
    start_condition1 = (temp['CURRENT DENSITY'] - temp['CURRENT DENSITY'].shift()) > 0.001
    start_condition2 = (temp['CURRENT DENSITY']) < 0.02
    start_condition = start_condition1.shift(-1, fill_value=False) & start_condition2
    startpoints = temp.loc[start_condition, :].copy()
    startpoints['type'] = 0
    # print(startpoints)
    # put them together
    pol_times = pd.concat([startpoints, endpoints])
    pol_times.sort_values(by=['Date'], inplace=True)
    # get only the start and endpoints that border each other
    change_condition = (pol_times['type'] - pol_times['type'].shift()) == 1
    pol_times = pol_times.loc[(change_condition | change_condition.shift(-1, fill_value=False)), :]
    # print(pol_times)
    # check if monotonically increasing for each interval
    t_start = []
    t_end = []
    # for each interval
    for i in range(len(pol_times.loc[pol_times['type'] == 0, 'Date'])):
        # get the start and end time of the overall interval
        start = pol_times.iloc[i*2]['Date']
        end = pol_times.iloc[(i*2)+1]['Date']
        # sample the data at nominal frequency between start and end
        sample = temp[(temp['Date'] >= start) & (temp['Date'] <= end)]
        # if there is less than 18 points, then dont use this interval
    #     print(len(sample))
        if len(sample) < 18:
            continue
        else:
            # make the rolling window size 10% of interval length
            win = int(len(sample)/10)
            # sample the data based on this window size and take the mean
            sample_sparse = sample['CURRENT DENSITY'].rolling(win, min_periods=1, center=True, step = win).mean() 
            # while the data is monotonically increasing (by more than 0.01)
            idx = 1
            while (sample_sparse.values[idx] - sample_sparse.values[idx-1]) >-0.1:
                # update the counter
                idx += 1
                if idx >= len(sample_sparse):
                    break
            # update the end point based on the counter if the counter did not reach the full len
            if idx < len(sample_sparse):
                end = sample.loc[sample_sparse.index[idx-1], 'Date']
                # resample the data at nominal rate using the updated endpoint
                sample = temp[(temp['Date'] >= start) & (temp['Date'] <= end)]
            # if the resulting pol interval has at least 18 points, then update the relevant lists
    #         print(len(sample))
            if len(sample) > 18:
                t_start.append(start)
                t_end.append(end)
                
    # specify the run duration completed at the polcurve start
    T1 = t_start
    T2 = t_end
    N=[]
    for i in T1:
    #     print(Vcell_data[Vcell_data['Date'] == i][['Date','DURATION_HR']])
        d = str(round(Vcell_data[Vcell_data['Date'] == i]['DURATION_HR'].iloc[0].item(),1))+" hrs"
        N.append(d)
    print('configured pol curve start, end & N (duration completed at the polcurve start)')
    #print(T1,T2,N)    
    return T1, T2, N


# In[23]:

# func to plot cd and identified polcurves chart
def cd_and_polcurve(Vcell_data, T1, T2, image_folder_name):
    mask = []
    for i,j in zip(T1,T2):
        temp_t = Vcell_data.loc[(Vcell_data['Date'] >= i ) & (Vcell_data['Date'] <= j), ['Date','CURRENT DENSITY']]
        mask.append(temp_t)

    stacked_df = pd.DataFrame(columns=temp_t.columns)
    for temp_t in mask:
        stacked_df = pd.concat([stacked_df, temp_t], ignore_index=True)
        #stacked_df.append(temp_t, ignore_index=True)
    stacked_df.rename(columns={'CURRENT DENSITY':'pol_curves'}, inplace = True) 


    Current_Density = go.Scattergl(
                           x = Vcell_data['Date'],
                           y = Vcell_data['CURRENT DENSITY'],
                           mode = 'markers',
                           marker=dict(color='#2ca02c'),
    #                        mode = 'lines',
                           name = 'Current Density (A/c'+'m\u00b2)',
                           yaxis = 'y1'
                           #text = df.column (this is the info displayed on hover)
                           )
    pol_curves = go.Scattergl(
                           x = stacked_df['Date'],
                           y = stacked_df['pol_curves'],
                           mode = 'markers',
    #                        marker=dict(color='#2ca02c'),
    #                        mode = 'lines',
                           name = 'pol_curves',
                           yaxis = 'y1'
                           #text = df.column (this is the info displayed on hover)
                           )

    # update the data list
    data = [Current_Density,pol_curves]

    # configure the layout
    layout = go.Layout(
                title = 'Current Density',
                #to define axis, can either use go.layout.etc or dict(etc)
                xaxis = go.layout.XAxis(
    #                                     title = 'Date',
                                        tickformat = "%d-%b"+'<br>'+" %H:00",
                                        #hoverformat = "%d/%m/%y"+'<br>'+"  %H:%M:%S",
                                        hoverformat = "%d-%b %H:00",
                                        showgrid = False,
                                        #range=[100,steady_state['Duration(hr)'].max()]
                                        ),    
                yaxis= dict(title='Current Density'+'<br>'+'(A/c'+'m\u00b2)',
    #                                  overlaying='y',
                                     side = 'left',
                                     showgrid = False,
    #                                  range = [0,6],
                                     rangemode='tozero'),
                legend= dict(xanchor="right",x=1.5,yanchor="top", y =2, bgcolor = 'rgba(0,0,0,0)'),
                template = 'simple_white',
                font = dict(size = 18)
                )


    # create the figure
    fig0 = go.Figure(data = data, layout = layout)
    fig0.update_xaxes(tickangle=90)
    
    fig0.update_layout(xaxis = dict(titlefont = dict(size = 18), tickfont = dict(size = 18)),
           yaxis = dict(titlefont = dict(size = 18), tickfont = dict(size = 18)),
           yaxis2 = dict(titlefont = dict(size = 18), tickfont = dict(size = 18)) 
                     )

    fig0.write_image('{}/Cd_&_polcurves.png'.format(image_folder_name), engine="orca")
    fig0.write_html('{}/Cd_&_polcurves.html'.format(image_folder_name), include_plotlyjs=('cdn'))
    # fig0.show()
    print('exported cd_&_identified_polcurves chart')


# In[24]:

# func to plot avg_stack_polcurve
def avg_stack_v_polcurve(N, T1, T2, Vcell_data, image_folder_name):
    n=N
    t1=T1
    t2=T2

    #creating lists to later make Min/Max Cell table from these
    N_temp=N.copy()
    T1_temp=T1.copy()
    T2_temp=T2.copy()

    fig = go.Figure()
    
    for n,t1,t2 in zip(n,t1,t2):

        pol_duration_df = Vcell_data[(Vcell_data.Date >= t1) & (Vcell_data.Date <= t2)]
        poldata_list = Stack_voltage_list.copy()
        poldata_list.extend(['Date','CURRENT SP','STACK CURRENT','CURRENT DENSITY'])
        pol_data = pol_duration_df[poldata_list]

        ## choose indexes that have only one row for a current SP  
        indexes = pol_data[pol_data.groupby("CURRENT SP")["CURRENT SP"].transform('size') == 1].index
        ## and then drop those rows
        pol_data.drop(indexes, axis=0, inplace=True)
        #df_polcurve1 = pd.DataFrame(columns = list(pol_data.columns))
        list_curr_sp = pol_data['CURRENT SP'].unique()
        df_temp_list = [pol_data[pol_data['CURRENT SP']==i].iloc[-2] for i in list_curr_sp]
        #df_temp_list = df_polcurve1 + df_temp_list
        #df_polcurve1 = pd.concat(df_temp_list)
        df_polcurve1 = pd.DataFrame(df_temp_list)
        #for i in list_curr_sp:
            #df_temp = pol_data[pol_data['CURRENT SP']==i]
            #df_polcurve1 = df_polcurve1.append(df_temp.iloc[-2])    
        df_polcurve1 = df_polcurve1.reset_index(drop=True)    
        ############################# added logic to remove decreasing CD
        ##create a list out of sc column
        polcurve1_SC = df_polcurve1['STACK CURRENT']
        polcurve1_SC = list(polcurve1_SC)
        # check if next element of list is decreasing , remove it
        index = 0
        while index < (len(polcurve1_SC)-1):
            diff = polcurve1_SC[index + 1] - polcurve1_SC[index]
            if diff < 3:
                polcurve1_SC.pop(index + 1)  # Remove the next element
            else:
                index += 1  # Move to the next element if no removal was made

        ## writing a condition that if unique data pnts in the polcurve data is less than 3, then drop that polcurvre and continue to next polcurve        
        if len(polcurve1_SC)<3:
            N_temp.remove(n)
            T1_temp.remove(t1)
            T2_temp.remove(t2)
            continue
        # else go for making the polcurve further    
        else:                
            # create a new dataframe out of polcurve1_SC list and join with the original df to pick required data rows.
            df_polcurve2 = pd.DataFrame()
            df_polcurve2['STACK CURRENT'] =  polcurve1_SC
            ## joining two dfs
            df_polcurve1 = df_polcurve2.merge(df_polcurve1, on='STACK CURRENT', how='left')
            ###############################
        #     print(df_polcurve1)
    #         df_polcurve1['Avg_V'] = df_polcurve1[Stack_voltage_list].sum(axis=1)/len(Stack_voltage_list)
            df_polcurve1['Avg_V'] = df_polcurve1[Stack_voltage_list].mean(axis=1)
            df_polcurve1['Max_V'] = df_polcurve1[Stack_voltage_list].max(axis=1)
            df_polcurve1['Min_V'] = df_polcurve1[Stack_voltage_list].min(axis=1)
    #         df_polcurve1.to_csv('df_polcurve1.csv')
            fig.add_trace(go.Scatter(
                                 x = df_polcurve1['CURRENT DENSITY'],
                                 y = df_polcurve1['Avg_V'],
                                 showlegend = True,
                                 mode = 'lines',
                                 name = n,
                                 yaxis = 'y1',

                                 error_y = dict(
                                     type = 'data',
                                     symmetric=False,
                                     visible = True,
                                     array = df_polcurve1['Max_V'] - df_polcurve1['Avg_V'],
                                     arrayminus = df_polcurve1['Avg_V'] - df_polcurve1['Min_V']
                                     ),
                                 customdata = np.stack((df_polcurve1['Min_V'], df_polcurve1['Max_V']), axis=-1),
        #                          customdata2= np.stack(df_polcurve1['Date']), 
                                 hovertemplate ='Max: %{customdata[1]:.2f}'+'Avg: %{y:.2f}'+'<br>'+'Min: %{customdata[0]:.2f}'
        #         +'<br>'+'Date: %{customdata2[0]:dd/mm/yyyy hh:mm:ss'
                                ))
            fig.update_xaxes(title = 'Current Density (A/c'+'m\u00b2)')
            fig.update_yaxes(title = 'Voltage (V)')

            fig.update_layout(title=dict(text="Pol. Curve (Avg of all running stacks)", font=dict(size=20)),
                              xaxis = dict(titlefont = dict(size = 18), tickfont = dict(size = 18)),
                              yaxis = dict(titlefont = dict(size = 18), tickfont = dict(size = 18))  
                              # temorary setup to prevent negative values in polcurve but figure it out later
                              #,yaxis_range = [-1.4,None])
                              )
#     fig.show()
    # fig.write_image('{}/ALL_Pol_Curve_Avg_all_stacks.png'.format(image_folder_name), engine="orca")
    fig.write_html('{}/ALL_Pol_Curve_Avg_all_stacks.html'.format(image_folder_name), include_plotlyjs=('cdn'))

    # make 3 selected polcurve plot
    first_trace_index = 0
    last_trace_index = len(fig.data)-1
    middle_trace_index = len(fig.data)//2
    selected_traces = []
    for i, trace in enumerate(fig.data):
        if i == first_trace_index or i == middle_trace_index or i == last_trace_index:
            selected_traces.append(trace)
    fig_selected_pc = go.Figure(data=selected_traces)
    fig_selected_pc.layout = fig.layout
#     fig_selected_pc.show()
    fig_selected_pc.write_image('{}/Pol_Curve_Avg_all_stacks.png'.format(image_folder_name), engine="orca")
    fig_selected_pc.write_html('{}/Pol_Curve_Avg_all_stacks.html'.format(image_folder_name), include_plotlyjs=('cdn'))
    print('exported Avg voltage polcurve plot')
    return N_temp, T1_temp, T2_temp


# In[25]:

# func to create min/max voltage cell table image for Avg_voltage_polcurve_plot
def MinMax_v_cell(N_temp, T1_temp, T2_temp, Vcell_data, image_folder_name):
    n = [N_temp[0], N_temp[len(N_temp)//2], N_temp[len(N_temp)-1]]
    t1=[T1_temp[0], T1_temp[len(N_temp)//2], T1_temp[len(N_temp)-1]]
    t2=[T2_temp[0], T2_temp[len(N_temp)//2], T2_temp[len(N_temp)-1]]
    
    MinMax_cell_df = pd.DataFrame()
    Pol_curve_duration = []
    Max_Current_density = []
    Max_V_cell = []
    Min_V_cell = []
    
    for n,t1,t2 in zip(n,t1,t2):
        pol_duration_df = Vcell_data[(Vcell_data.Date >= t1) & (Vcell_data.Date <= t2)]
        
        Pol_curve_duration.append(n)
        Max_Current_density.append(round(max(pol_duration_df['CURRENT DENSITY']),2))

        pol_df = Vcell_data[Vcell_data['Date'] == t1][Stack_voltage_list]

        max_v_cell = pol_df.idxmax(axis=1).values
        max_v_cell_str = ", ".join(str(num) for num in max_v_cell)

        min_v_cell = pol_df.idxmin(axis=1).values
        min_v_cell_str = ", ".join(str(num) for num in min_v_cell)

        Max_V_cell.append(max_v_cell_str)
        Min_V_cell.append(min_v_cell_str)

    MinMax_cell_df['Pol curve duration'] = Pol_curve_duration
    MinMax_cell_df['Max Current density (A/cm2)'] = Max_Current_density
    MinMax_cell_df['Max_V_cell'] = Max_V_cell
    MinMax_cell_df['Min_V_cell'] = Min_V_cell

    # MinMax_cell_df

    ######### converting the table dataframe in a cropped image
    header_values = ['<b>' + col + '</b>' for col in MinMax_cell_df.columns]

    fig = go.Figure(data=[go.Table(
        header=dict(values = header_values,
    #                 fill_color='paleturquoise',
                    fill_color='lightgray',
                    align='center',
                    line=dict(color='black', width=1),
                    font=dict(color='black', size=12)
                   ), 
        cells=dict(values=[MinMax_cell_df['Pol curve duration'],
                        MinMax_cell_df['Max Current density (A/cm2)'], MinMax_cell_df['Max_V_cell'], MinMax_cell_df['Min_V_cell']],
                   fill_color='lavender',
                   align='center',
                   line=dict(color='black', width=1)))
                         ])

    fig.update_layout(
        margin=dict(t=0, b=0, l=0, r=0),
        xaxis=dict(visible=False),
        yaxis=dict(visible=False),
        showlegend=False
    )

    # fig.write_image('./images/MinMax_cell.png')
    fig.write_image('{}/MinMax_cell.png'.format(image_folder_name), engine="orca")

    # Open the image
#     image = Image.open("./"+image_folder_name + "/"+"MinMax_cell.png")
    image = Image.open(image_folder_name + "/MinMax_cell.png")
    # Define the crop region (left, upper, right, lower)
    crop_region = (0, 0, 700, 108)
    
    # Crop the image
    MinMax_cell = image.crop(crop_region)
    filename = "MinMax_cell.png"
#     folder_path = './'+ image_folder_name + '/'
    folder_path = image_folder_name + '/'

    
    MinMax_cell.save(folder_path + filename)
    print('MinMax_V_cell_table image exported')


# In[26]:

# func to create stack wise pol curve
def stack_wise_polcurve(N, T1, T2, Vcell_data, image_folder_name):
    cell_range = {'L1':Stack_voltage_list}
    running_cell_positions = {}
    for position, cell in cell_range.items():
        #check if all the cols for a stack is null completely, then pass(do not publish)   
        is_cols_blank = Vcell_data[cell].isnull().all()
        if is_cols_blank.all():
            pass
        else:
            running_cell_positions[position] = cell        
# print(str(running_cell_positions))
    for m,i in running_cell_positions.items():    
            n = N 
            t1 = T1
            t2 = T2
            fig = go.Figure()
            for n,t1,t2 in zip(n,t1,t2):
                pol_duration_df = Vcell_data[(Vcell_data.Date >= t1) & (Vcell_data.Date <= t2)]
                poldata_list = i.copy()
                poldata_list.extend(['Date','CURRENT SP','STACK CURRENT','CURRENT DENSITY'])
                pol_data = pol_duration_df[poldata_list]
                indexes = pol_data[pol_data.groupby("CURRENT SP")["CURRENT SP"].transform('size') == 1].index
                pol_data.drop(indexes, axis=0, inplace=True)
            
                #df_polcurve1 = pd.DataFrame(columns = list(pol_data.columns))
                list_curr_sp = pol_data['CURRENT SP'].unique()
                
                df_temp_list = [pol_data[pol_data['CURRENT SP']==i].iloc[-2] for i in list_curr_sp]
                df_polcurve1  = pd.DataFrame(df_temp_list)
                #for j in list_curr_sp:
                 #   df_temp = pol_data[pol_data['CURRENT SP']==j]
                 #   df_polcurve1 = df_polcurve1.append(df_temp.iloc[-2])
                df_polcurve1 = df_polcurve1.reset_index(drop=True)    
    #             print(df_polcurve1)   
                ################################## added logic to remove decreasing CD
                ##create a list out of sc column
                polcurve1_SC = df_polcurve1['STACK CURRENT']
                polcurve1_SC = list(polcurve1_SC)
                # check if next element of list is decreasing , remove it
                index = 0
                while index < (len(polcurve1_SC)-1):
                    diff = polcurve1_SC[index + 1] - polcurve1_SC[index]
                    if diff < 3:
                        polcurve1_SC.pop(index + 1)  # Remove the next element
                    else:
                        index += 1  # Move to the next element if no removal was made
                ## writing a condition that if unique data pnts in the polcurve data is less than 3, then drop that polcurvre and continue to next polcurve        
                if len(polcurve1_SC)<3:
    #                 N.remove(n)
    #                 T1.remove(t1)
    #                 T2.remove(t2)
                    continue
                # else go for making the polcurve further    
                else:               
                    # create a new dataframe out of polcurve1_SC list and join with the original df to pick required data rows.
                    df_polcurve2 = pd.DataFrame()
                    df_polcurve2['STACK CURRENT'] =  polcurve1_SC
                    ## joining two dfs
                    df_polcurve1 = df_polcurve2.merge(df_polcurve1, on='STACK CURRENT', how='left')
                    ##################################    

                    df_polcurve1['Avg_V'] = df_polcurve1[i].mean(axis=1)
                    df_polcurve1['Max_V'] = df_polcurve1[i].max(axis=1)
                    df_polcurve1['Min_V'] = df_polcurve1[i].min(axis=1)
    #                 df_polcurve1.to_csv(m+'_polcurve.csv')
                   # if there is no data (may be stacks removed), then do not plot the trace or legend
                    if df_polcurve1['Avg_V'].any():
                        fig.add_trace(go.Scatter(
                                             x = df_polcurve1['CURRENT DENSITY'],
                                             y = df_polcurve1['Avg_V'],
                                             mode = 'lines',
                                             showlegend = True,
                                             name = n,
                                             yaxis = 'y1',
                                             error_y = dict(
                                                 type = 'data',
                                                 symmetric=False,
                                                 visible = True,
                                                 array = df_polcurve1['Max_V'] - df_polcurve1['Avg_V'],
                                                 arrayminus = df_polcurve1['Avg_V'] - df_polcurve1['Min_V']
                                                 ),

                                             customdata = np.stack((df_polcurve1['Min_V'], df_polcurve1['Max_V']), axis=-1),
                                             hovertemplate ='Max: %{customdata[1]:.2f}'+'Avg: %{y:.2f}'+'<br>'+'Min: %{customdata[0]:.2f}'
                                            ))
                        fig.update_xaxes(title = 'Current Density (A/c'+'m\u00b2)')
                        fig.update_yaxes(title = 'Voltage (V)')

                        fig.update_layout(title=dict(text=m+" Pol Curve", font=dict(size=20)),
                                          xaxis = dict(titlefont = dict(size = 18), tickfont = dict(size = 18)),
                                          yaxis = dict(titlefont = dict(size = 18), tickfont = dict(size = 18)))
                    else:
                        fig.add_trace(go.Scatter(
                                             x = df_polcurve1['CURRENT DENSITY'],
                                             y = df_polcurve1['Avg_V'],
                                             mode = 'lines',
                                             showlegend = False))

    #             fig.show()
    #         fig.write_image('{}/{}_ALL_pol_curve.png'.format(image_folder_name, m), engine="orca")
            fig.write_html('{}/{}_ALL_pol_curve.html'.format(image_folder_name, m), include_plotlyjs=('cdn'))

            # make 3 selected polcurve plot
            first_trace_index = 0
            last_trace_index = len(fig.data)-1
            middle_trace_index = len(fig.data)//2
            selected_traces = []
            for i, trace in enumerate(fig.data):
                if i == first_trace_index or i == middle_trace_index or i == last_trace_index:
                    selected_traces.append(trace)
            fig_selected_pc = go.Figure(data=selected_traces)
            fig_selected_pc.layout = fig.layout
#             fig_selected_pc.show()
            fig_selected_pc.write_image('{}/{}_pol_curve.png'.format(image_folder_name,m), engine="orca")
            fig_selected_pc.write_html('{}/{}_pol_curve.html'.format(image_folder_name,m), include_plotlyjs=('cdn')) 
    print('exported stack wise polcurve charts')
    return running_cell_positions


# In[27]:

# func to create SSE chart
def stack_specific_energy(Vcell_data, image_folder_name):
    # No_of_Stacks = len(Stack_voltage_list)    
    Vcell_data['SUM_STACK_VOLTAGE'] = Vcell_data[Stack_voltage_list].sum(axis=1)
    # print(sys_data[sys_data['Date'] == '05-11-2023 12:26:40'][['Date','SUM_STACK_VOLTAGE']])

    Vcell_data['no_of_stacks'] = Vcell_data[Stack_voltage_list].apply(lambda row: row[row != 0].count(), axis=1)

    ## if no of stacks is not a multiple of 56 , then drop that row (based on the fact that there will be always 56 stacks present in the system)
    #Vcell_data  = Vcell_data[Vcell_data['no_of_stacks']%56 == 0]

    # sys_data['STACK SPECIFIC ENERGY'] = sys_data['SUM_STACK_VOLTAGE']/(0.04*No_of_Stacks) [this was old calculation which was corrected later]
    Vcell_data['STACK SPECIFIC ENERGY'] = Vcell_data['SUM_STACK_VOLTAGE']/(0.0376*Vcell_data['no_of_stacks'])
    # print(sys_data[sys_data['Date'] == '05-11-2023 12:26:40'][['Date','STACK SPECIFIC ENERGY']])
    
    ##avg_per_hour = Vcell_data.groupby(['date_&_hrs'])[['CURRENT DENSITY','STACK SPECIFIC ENERGY','order']].mean().reset_index()    
    #avg_per_hour = Vcell_data.loc[:,['Date','CURRENT DENSITY','STACK SPECIFIC ENERGY','order']].resample('H', on='Date').mean().reset_index()
    
    Stack_Specific_energy = go.Scattergl(
                           x = Vcell_data['Date'],
                           y = Vcell_data['STACK SPECIFIC ENERGY'],
                           mode = 'markers',
                           #marker=dict(size=3),
                           #mode = 'lines',
                           name = 'Stack Specific Energy (kWh/kg)',
                           yaxis = 'y1'
                           #text = df.column (this is the info displayed on hover)
                           )
    Current_Density = go.Scattergl(
                           ##x = avg_per_hour['date_&_hrs'],
                           x = Vcell_data['Date'],
                           y = Vcell_data['CURRENT DENSITY'],
                           mode = 'markers',
    #                        mode = 'lines',
                           name = 'Current Density (A/c'+'m\u00b2)',
                           yaxis = 'y2'
                           #text = df.column (this is the info displayed on hover)
                           )

    # update the data list
    data = [Stack_Specific_energy, Current_Density]

    # configure the layout
    layout = go.Layout(
                title = 'LHC Stack Specific'+'<br>'+'Energy (DC)',
                #to define axis, can either use go.layout.etc or dict(etc)
                xaxis = go.layout.XAxis(
    #                                     title = 'Date',
                                        nticks= 4,
                                        tickformat = "%d-%b"+'<br>'+" %H:00",
                                        hoverformat = "%d-%b %H:00",
                                        showgrid = False,
                                        #range=[100,steady_state['Duration(hr)'].max()]
                                        ),    
                yaxis = dict(title='Stack Specific'+'<br>'+'Energy (kWh/kg)',
                                    showgrid = False,
                                    rangemode='tozero'),
                yaxis2= dict(title='Current Density'+'<br>' +'(A/c'+'m\u00b2)',
                                     overlaying='y',
                                     side = 'right',
                                     range = [0,2],
                                     showgrid = False,
                                     rangemode='tozero'),
    #                                  range=[0,600]),
                legend = dict(xanchor="right",x=1.4,yanchor="top", y =1.3, bgcolor = 'rgba(0,0,0,0)'),
                template = 'simple_white',
                font = dict(size = 18)
                )

    # create the figure
    fig4 = go.Figure(data = data, layout = layout)

    fig4.update_layout(xaxis = dict(titlefont = dict(size = 18), tickfont = dict(size = 18),
                                    tickformat = "%d-%b"+'<br>'+" %H:00",
                                    hoverformat = "%d-%b %H:00"),
                      yaxis = dict(titlefont = dict(size = 18), tickfont = dict(size = 18)),
                      yaxis2 = dict(titlefont = dict(size = 18), tickfont = dict(size = 18)),
                      )
    
    fig4.write_image('{}/LHC_Stack_Specific_Energy.png'.format(image_folder_name), engine="orca")
    fig4.write_html('{}/LHC_Stack_Specific_Energy.html'.format(image_folder_name), include_plotlyjs=('cdn'))
    # fig4.show()
    print('exported stack specific energy chart')


# In[28]:
def stack_specific_energy_new(sys_data, image_folder_name):
    # No_of_Stacks = len(Stack_voltage_list)    
    # columns_to_check = sys_data.loc[:,'C1-1','C16-16'].columns
    Vcell_data['SUM_STACK_VOLTAGE'] = Vcell_data[Stack_voltage_list].sum(axis=1)
    # print(sys_data[sys_data['Date'] == '05-11-2023 12:26:40'][['Date','SUM_STACK_VOLTAGE']])

    sys_data['no_of_stacks'] = sys_data[Stack_voltage_list].apply(lambda row: row[row != 0].count(), axis=1)

    ## if no of stacks is not a multiple of 56 , then drop that row (based on the fact that there will be always 56 stacks present in the system)
    #sys_data  = sys_data[sys_data['no_of_stacks']%56 == 0]

    # sys_data['STACK SPECIFIC ENERGY'] = sys_data['STACK VOLTAGE']/(0.04*No_of_Stacks) [this was old calculation which was corrected later]
    sys_data['STACK SPECIFIC ENERGY'] = sys_data['SUM_STACK_VOLTAGE']/(0.0376*sys_data['no_of_stacks'])
    # print(sys_data[sys_data['Date'] == '05-11-2023 12:26:40'][['Date','STACK SPECIFIC ENERGY']])
    
    ##avg_per_hour = Vcell_data.groupby(['date_&_hrs'])[['CURRENT DENSITY','STACK SPECIFIC ENERGY','order']].mean().reset_index()    
    #avg_per_hour = sys_data.loc[:,['Date','CURRENT DENSITY','STACK SPECIFIC ENERGY','order']].resample('H', on='Date').mean().reset_index()
    
    Stack_Specific_energy = go.Scattergl(
                           x = sys_data['Date'],
                           y = sys_data['STACK SPECIFIC ENERGY'],
                           mode = 'markers',
                           marker=dict(size=2),
                           #mode = 'lines',
                           name = 'SSE',
                           yaxis = 'y1'
                           #text = df.column (this is the info displayed on hover)
                           )
    Current_Density = go.Scattergl(
                           ##x = avg_per_hour['date_&_hrs'],
                           x = sys_data['Date'],
                           y = sys_data['CURRENT DENSITY'],
                           mode = 'markers',
                           marker=dict(size=2),
    #                        mode = 'lines',
                           name = 'Cd',
                           yaxis = 'y2'
                           #text = df.column (this is the info displayed on hover)
                           )

    # update the data list
    data = [Stack_Specific_energy, Current_Density]

    # configure the layout
    layout = go.Layout(
                title = 'LHC Stack Specific Energy (DC)',
                #to define axis, can either use go.layout.etc or dict(etc)
                xaxis = go.layout.XAxis(
    #                                     title = 'Date',
                                        nticks= 5,
                                        tickformat = "%d-%b"+'<br>'+" %H:00",
                                        hoverformat = "%d-%b %H:00",
                                        showgrid = False,
                                        #range=[100,steady_state['Duration(hr)'].max()]
                                        ),    
                yaxis = dict(title='Stack Specific Energy (kWh/kg)',
                                    showgrid = False,
                                    rangemode='tozero'),
                yaxis2= dict(title='Cd (A/c'+'m\u00b2)',
                                     overlaying='y',
                                     side = 'right',
                                     range = [0,2],
                                     showgrid = False,
                                     rangemode='tozero'),
    #                                  range=[0,600]),
                legend = dict(xanchor="left",x=0,yanchor="top", y =1.1, bgcolor = 'rgba(0,0,0,0)',orientation='h'),
                template = 'simple_white',
                font = dict(size = 18),
                width=850
                )

    # create the figure
    fig4 = go.Figure(data = data, layout = layout)

    fig4.update_layout(xaxis = dict(titlefont = dict(size = 18), tickfont = dict(size = 18)),
                                    
                      yaxis = dict(titlefont = dict(size = 18), tickfont = dict(size = 18)),
                      yaxis2 = dict(titlefont = dict(size = 18), tickfont = dict(size = 18)),
                      )
    
    fig4.write_image('{}/LHC_Stack_Specific_Energy_new.png'.format(image_folder_name), engine="orca")
    fig4.write_html('{}/LHC_Stack_Specific_Energy_new.html'.format(image_folder_name), include_plotlyjs=('cdn'))
    # fig4.show()
    print('exported stack specific energy chart')
#%%
def avg_min_max_stack_v_old_singleaxis(sys_data, image_folder_name):

    avg_per_hour = sys_data.loc[:,['Date','CURRENT DENSITY','Avg_V_C1_1_C8_16', 'Max_V_C1_1_C8_16', 'Min_V_C1_1_C8_16','std_dev']].resample('H', on='Date').mean().reset_index()
    # color_dict = {'Avg':fig.layout['template']['layout']['colorway'][0] ,
    #               'Max':fig.layout['template']['layout']['colorway'][1] ,
    #               'Min':fig.layout['template']['layout']['colorway'][2] ,
    #           'std_dev':fig.layout['template']['layout']['colorway'][3] ,
    #               'i' :fig.layout['template']['layout']['colorway'][4] } 

    # defining current density
    Current_Density = go.Scattergl(
                            x = sys_data['Date'],
                            y = sys_data['CURRENT DENSITY'],
                            mode = 'markers',
    #                        mode = 'lines',
                            name = 'Cd',
                            marker=dict(size=3,color='black'),
                            yaxis='y2'
                            #text = df.column (this is the info displayed on hover)
                            )    
    # defining std. dev
    std_dev_trace = go.Scattergl(
                            x = sys_data.Date,
                            y = sys_data['std_dev'],
                            mode = 'markers',
                            name = 'Std Dev',
                            #line = dict(color = 'rgb(255, 127, 14)', shape='hv'),
                            marker=dict(size=3,color='gray'),
                            yaxis='y2'
                            )
    Max_StackL1 = go.Scattergl(
                            x = sys_data['Date'],
                            y = sys_data['Max_V_C1_1_C8_16'],
                            mode = 'markers',
                            #mode = 'lines',
                            name = 'Max V',
                            #legendgroup=3,
                            marker=dict(size=3),
                            yaxis = 'y1'
                            #text = df.column (this is the info displayed on hover)
                            )   
    Min_StackL1 = go.Scattergl(
                            x = sys_data['Date'],
                            y = sys_data['Min_V_C1_1_C8_16'],
                            mode = 'markers',
                            #mode = 'lines',
                            name = 'Min V',
                            #legendgroup=3,
                            #showlegend=False,
                            marker=dict(size=3),
                            yaxis = 'y1'
                            #text = df.column (this is the info displayed on hover)
                            )    
    Avg_StackL1 = go.Scattergl(
                            x = sys_data['Date'],
                            y = sys_data['Avg_V_C1_1_C8_16'],
                            mode = 'markers',
                            #mode = 'lines',
                            #legendgroup=3,
                            name = 'Avg V',
                            marker=dict(size=3),
                            yaxis = 'y1'
                            #text = df.column (this is the info displayed on hover)
                            )
    dashed_line = {
            'type': 'line',
            'xref': 'paper',
            'yref': 'y',
            'x0': 0,
            'y0': 2,
            'x1': 1,
            'y1': 2,
            'line': {
                'color': 'black',
                'width': 2,
                'dash': 'dash',
            },
        }
    
    data_plot = [Current_Density,
                 std_dev_trace,
                 Max_StackL1, Min_StackL1, Avg_StackL1]
    
    # configure the layout
    layout = go.Layout(
                #title = 'LHC Pressure & Conductivity chart',
                #title_font=dict(weight='bold'),
                #to define axis, can either use go.layout.etc or dict(etc)
                xaxis = go.layout.XAxis(
    #                                     title = 'Date',
                                        nticks= 5,
                                        tickformat = "%d-%b"+'<br>'+" %H:00",
                                        #hoverformat = "%d/%m/%y"+'<br>'+"  %H:%M:%S",
                                        hoverformat = "%d-%b %H:00",
                                        showgrid = False,
                                        #tickfont=dict(size=16, color='black')  # Adjust font properties as needed
                                        #tickangle=90
                                        #range=[100,steady_state['Duration(hr)'].max()]
                                        ),
                yaxis = dict(title='Voltage (V)',
    #                                 font = 15,
                                    showgrid = False,
                                    rangemode='tozero',
                                    range=[1.4, 2]),
                yaxis2= dict(title='Std Dev (V) & Cd'+' (A/c'+'m\u00b2)',
    #                                  overlaying='y',
                                     side = 'right',
                                     showgrid = False,
                                     rangemode='tozero',
    #                                  range=[0,12.5]
                            ),
              #legend= dict(xanchor="right",x=1.3, yanchor="top", y =1.5, bgcolor = 'rgba(0,0,0,0)'),
              legend= dict(xanchor="right",x=1.35, yanchor="top", y =0.8, bgcolor = 'rgba(0,0,0,0)'),
              template = 'simple_white',
              font = dict(size = 18),
              shapes=[dashed_line],
              width=850
                  )
    # create the figure
    fig = go.Figure(data = data_plot, layout = layout)

    fig.update_layout(xaxis = dict(titlefont = dict(size = 18), tickfont = dict(size = 18)),
            yaxis = dict(titlefont = dict(size = 18), tickfont = dict(size = 18)),
            yaxis2 = dict(titlefont = dict(size = 18), tickfont = dict(size = 18)))
    
    fig.write_image('{}/LHC_Stack_Voltage_&_Current_Density.png'.format(image_folder_name), engine="orca")
    fig.write_html('{}/LHC_Stack_Voltage_&_Current_Density.html'.format(image_folder_name), include_plotlyjs=('cdn'))
    # fig7.show()
    print('exported LHC_Stack_Voltage chart')
#%%
def avg_min_max_stack_v_Ben(Vcell_data, sys_data, image_folder_name):

    #avg_per_hour = Vcell_data.loc[:,['Date','CURRENT DENSITY','Avg_V_C9_1_C16_16', 'Max_V_C9_1_C16_16', 'Min_V_C9_1_C16_16']].resample('H', on='Date').mean().reset_index()
        
    # define a new figure
    fig = sp.make_subplots(rows=2, cols=1, shared_xaxes=True, vertical_spacing=0.03, row_heights=[2/3, 1/3])
    fig.update_layout(
                #title = 'LHC Cell Voltage & Current Density',
                #to define axis, can either use go.layout.etc or dict(etc)
                
                xaxis4 = dict(#title = 'Date',
                              nticks= 4,
                              tickformat = "%d-%b"+'<br>'+" %H:00",
                              #hoverformat = "%d/%m/%y"+'<br>'+"  %H:%M:%S",
                              hoverformat = "%d-%b %H:00",
                              showgrid = False),                                                                                
                yaxis = dict(title='Voltage (V)',
    #                                font = 15,
                                    range = [1.4, 2],
                                    showgrid = False                        
                                    #rangemode='tozero'
                                    ),
                yaxis2 = dict(title='Std Dev (V) & Cd (A/c'+'m\u00b2)',
    #                                 font = 15,
                                    #range = [1.3, 2.3],
                                    showgrid = False
                                    #rangemode='tozero'
                                    ),

                #legend= dict(xanchor="right",x=1.5,yanchor="top", y =2, bgcolor = 'rgba(0,0,0,0)'),               
                #legend= dict(xanchor="right",yanchor="top", y=1.38, bgcolor = 'rgba(0,0,0,0)'),
                legend= dict(bgcolor = 'rgba(0,0,0,0)'),
                template = 'simple_white',
                font = dict(size = 12)
                )

    fig.update_layout(xaxis = dict(titlefont = dict(size = 12), tickfont = dict(size = 12)),
           yaxis = dict(titlefont = dict(size = 12), tickfont = dict(size = 12)),
           yaxis2 = dict(titlefont = dict(size = 12), tickfont = dict(size = 12)) )
    # color_dict = {'U1':fig.layout['template']['layout']['colorway'][0] ,
    #               'U2':fig.layout['template']['layout']['colorway'][1] ,
    #               'L2':fig.layout['template']['layout']['colorway'][2] ,
    #               'L1':fig.layout['template']['layout']['colorway'][3] ,
    #               'i' :fig.layout['template']['layout']['colorway'][4] } 
    
    # defining max voltage traces
    
    Max_StackL1 = go.Scattergl(
                            x = Vcell_data['Date'],
                            y = Vcell_data['Max_V'],
                            mode = 'markers',
                            #mode = 'lines',
                            name = 'Max V',
                            #legendgroup=3,
                            marker=dict(size=3.5),
                            yaxis = 'y1'
                            #text = df.column (this is the info displayed on hover)
                            )   
    #adding traces to the figure
    fig.add_trace(Max_StackL1, row=1, col=1)
    
    # defining avg voltage traces
    
    Avg_StackL1 = go.Scattergl(
                            x = Vcell_data['Date'],
                            y = Vcell_data['Avg_V'],
                            mode = 'markers',
                            #mode = 'lines',
                            #legendgroup=3,
                            name = 'Avg V',
                            marker=dict(size=3.5),
                            yaxis = 'y1'
                            #text = df.column (this is the info displayed on hover)
                            )

    #adding traces to the figure
    fig.add_trace(Avg_StackL1, row=1, col=1)

    # defining Min voltage traces 
    Min_StackL1 = go.Scattergl(
                             x = Vcell_data['Date'],
                             y = Vcell_data['Min_V'],
                             mode = 'markers',
                             #mode = 'lines',
                             name = 'Min V',
                             #legendgroup=3,
                             #showlegend=False,
                             marker=dict(size=3.5, color='#9467bd'),
                             yaxis = 'y1'
                             #text = df.column (this is the info displayed on hover)
                             )   
     #adding traces to the figure
    fig.add_trace(Min_StackL1, row=1, col=1)    
    # defining current density
    Current_Density = go.Scattergl(
                           x = sys_data['Date'],
                           y = sys_data['CURRENT DENSITY'],
                           mode = 'markers',
    #                        mode = 'lines',
                           name = 'Cd',
                           marker=dict(color = 'rgb(214, 39, 40)',size=3.5),
                           yaxis='y2'
                           #text = df.column (this is the info displayed on hover)
                           )
    fig.add_trace(Current_Density, row=2, col=1)
    
    std_dev_trace = go.Scattergl(
                            x = sys_data.Date,
                            y = Vcell_data['std_dev'],
                            mode = 'markers',
                            name = 'Std Dev',
                            #line = dict(color = 'rgb(255, 127, 14)', shape='hv'),
                            marker=dict(size=3.5, color='rgb(68, 170, 153)'),
                            yaxis='y2'
                            )
    #adding cd trace to the figure
    fig.add_trace(std_dev_trace, row=2, col=1)
    
    fig.write_image('{}/LHC_Stack_Voltage_&_Current_Density.png'.format(image_folder_name), engine="orca")
    fig.write_html('{}/LHC_Stack_Voltage_&_Current_Density.html'.format(image_folder_name), include_plotlyjs=('cdn'))
    # fig7.show()
    print('exported LHC_Stack_Voltage chart Ben')


#%%
#  function for publishing charts & ppt
def publish_presentation(template_ppt_fn, image_folder_name, ppt_fn, output_dir, total_duration_hrs, Start_time, End_time, Stack_information, running_cell_positions):
    prs = Presentation('{}.pptx'.format(template_ppt_fn))
    
    ## edit first table (only table) of first slide       
    # Select the slide containing the table (e.g., slide 0)
    slide = prs.slides[0]
    # Select the table in the slide (e.g., table 0)
    table1 = slide.shapes[3].table
    today = date.today()
    # write body cells
    table1.cell(0, 0).text = "Date    "+str(today.strftime("%d/%m/%Y"))
    now = datetime.now()
    dt_string = now.strftime("%H:%M") + "hrs.(IST)"
    table1.cell(1, 0).text = "Time  "+dt_string

    tables_list = [table1]
    for j in tables_list:
        for row in j.rows:
            for cell in row.cells:

                # Set font size
                cell.text_frame.paragraphs[0].font.size = Pt(16)  # Change the font size to your desired value

                # Set font style
                cell.text_frame.paragraphs[0].font.name = 'Quicksand (Headings)'  # Change the font style to your desired font

                # Set alignment
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Change the alignment to your desired value

                # set color to black
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # (R, G, B) = (0, 0, 0) for black
                cell.text_frame.paragraphs[0].font.bold = True
       
    ## edit first table of second slide       
    slide = prs.slides[1]
    table2 = slide.shapes[0].table
    table2.cell(0, 6).text = str(int(total_duration_hrs))+" hrs"
    table2.cell(1, 5).text = Start_time
    table2.cell(1, 7).text = End_time
    
    for cell in [table2.cell(0, 6), table2.cell(1, 5), table2.cell(1,7)]:

        # Set font size
        cell.text_frame.paragraphs[0].font.size = Pt(16)  # Change the font size to your desired value

        # Set font style
        cell.text_frame.paragraphs[0].font.name = 'Quicksand (Headings)'  # Change the font style to your desired font

        # Set alignment
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Change the alignment to your desired value

        # set color to black
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # (R, G, B) = (0, 0, 0) for black
        cell.text_frame.paragraphs[0].font.bold = False
    
    table3 = slide.shapes[1].table
#     rows_to_modify = table3.rows[1:4]
    rows_to_modify = range(1,len(Stack_information)+1)
    # Access all columns except first
    columns_to_modify = range(1, len(Stack_information.columns))

    for row in rows_to_modify:
        for col_idx in columns_to_modify: 
            cell = table3.cell(row, col_idx)
            cell_content = Stack_information.iat[row-1, col_idx]
            if not pd.isna(cell_content):
#               print(type(cell_content))
                if type(cell_content) is datetime:
                    cell_content = cell_content.strftime('%d-%b')
                cell.text = cell_content
                # Set font size
                cell.text_frame.paragraphs[0].font.size = Pt(12)  # Change the font size to your desired value

                # Set font style
                cell.text_frame.paragraphs[0].font.name = 'Quicksand (Headings)'  # Change the font style to your desired font

                # Set alignment
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Change the alignment to your desired value

                # set color to black
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # (R, G, B) = (0, 0, 0) for black
                cell.text_frame.paragraphs[0].font.bold = False


    # defining 3rd slide            
    slide = prs.slides[2]
    height = Inches(3.51)
    width = Inches(6.47)
    row1_y = Inches(0.00)
    row2_y = Inches(3.46)
    column1_x = Inches(0.00)
    column2_x = Inches(6.87)
    #column1_x_1=Inches(0)
    # place the LHC_Stack_Specific_Energy Chart
    pic = slide.shapes.add_picture('{}/LHC_Stack_Specific_Energy_new.png'.format(image_folder_name), column1_x, row2_y, width=width, height=height)
    pic.click_action.hyperlink.address = '{}/LHC_Stack_Specific_Energy_new.html'.format(image_folder_name)

    # place the LHC_TEMP_&_Level_Sensor_Chart
    pic = slide.shapes.add_picture('{}/LHC_TEMP_&_Level_Sensor_Chart.png'.format(image_folder_name), column1_x, row1_y,width=width, height=height)
    pic.click_action.hyperlink.address = '{}/LHC_TEMP_&_Level_Sensor_Chart.html'.format(image_folder_name)

    # place the LHC_O2_&_H2_leak_Sensor_Chart
    pic = slide.shapes.add_picture('{}/LHC_H2_&_O2_leak.png'.format(image_folder_name), column2_x, row2_y,width=width, height=height)
    pic.click_action.hyperlink.address = '{}/LHC_H2_&_O2_leak.html'.format(image_folder_name)

    # place the LWC Pressure_&_Conductivity chart
    pic = slide.shapes.add_picture('{}/LHC_Pressure_&_Conductivity_Chart.png'.format(image_folder_name), column2_x, row1_y,width=width, height=height)
    pic.click_action.hyperlink.address = '{}/Steady_State_Chart.html'.format(image_folder_name)

    # defining 4th slide
    slide = prs.slides[3]
    height = Inches(6.96)
    width = Inches(10.1)
    row1_y = Inches(0.53)
    #row2_y = Inches(4.18)
    column1_x = Inches(0)
    #column2_x = Inches(7.15)

    # place the LHC_Stack_Voltage_chart
    pic = slide.shapes.add_picture('{}/LHC_Stack_Voltage_&_Current_Density.png'.format(image_folder_name), column1_x, row1_y,width=width, height=height)
    pic.click_action.hyperlink.address = '{}/LHC_Stack_Voltage_&_Current_Density.html'.format(image_folder_name)
    
    #pic = slide.shapes.add_picture('{}/LHC_Avg_Stack_Voltage_&_Current_Density.png'.format(image_folder_name), column2_x, row1_y,width=width, height=height)
    #pic.click_action.hyperlink.address = '{}/LHC_Avg_Stack_Voltage_&_Current_Density.html'.format(image_folder_name)

    # place the LHC_Max_Stack_Voltage chart
    #pic = slide.shapes.add_picture('{}/LHC_Max_Stack_Voltage_&_Current_Density.png'.format(image_folder_name), column1_x, row1_y,width=width, height=height)
    #pic.click_action.hyperlink.address = '{}/LHC_Max_Stack_Voltage_&_Current_Density.html'.format(image_folder_name)

    # place the LHC_Min_Stack_Voltage chart
    #pic = slide.shapes.add_picture('{}/LHC_Min_Stack_Voltage_&_Current_Density.png'.format(image_folder_name), column1_x, row2_y,width=width, height=height)
    #pic.click_action.hyperlink.address = '{}/LHC_Min_Stack_Voltage_&_Current_Density.html'.format(image_folder_name)

    # place the Current Density Chart
    #pic = slide.shapes.add_picture('{}/Current_Density.png'.format(image_folder_name), column2_x, row2_y,width=width, height=height)
    #pic.click_action.hyperlink.address = '{}/Current_Density.html'.format(image_folder_name)

    #definig 5th slide
    slide = prs.slides[4]
    height = Inches(7.02)
    width = Inches(9.51)
    row1_y = Inches(0.47)
    column1_x = Inches(0)
    column2_x = Inches(7.04)

    # place the Stack_ASR_Chart
    pic = slide.shapes.add_picture('{}/Stack_ASR_Chart.png'.format(image_folder_name), column1_x, row1_y,width=width, height=height)
    pic.click_action.hyperlink.address = '{}/Stack_ASR_Chart.html'.format(image_folder_name)

    # place the Stack_ASR(max)_Chart
    #pic = slide.shapes.add_picture('{}/Stack_ASR(max)_Chart.png'.format(image_folder_name), column2_x, row1_y,width=width, height=height)
    #pic.click_action.hyperlink.address = '{}/Stack_ASR(max)_Chart.html'.format(image_folder_name)

    # defining 6th slide
    slide = prs.slides[5]
    height = Inches(5.28)
    width = Inches(7.04)
    row1_y = Inches(0.5)
    column1_x = Inches(0.1)
    
    height_1 = Inches(1.5)
    width_1 = Inches(10.42)
    row1_y_1 = Inches(5.88)
    column1_x_1 = Inches(1.31)
    
    avg_pol_curve_file = '{}/Pol_Curve_Avg_all_stacks.png'.format(image_folder_name)
    minmax_cell_file = '{}/MinMax_cell.png'.format(image_folder_name)
    
    if os.path.exists(avg_pol_curve_file):
        pic = slide.shapes.add_picture(avg_pol_curve_file, column1_x, row1_y, width=width, height=height)

    if os.path.exists(minmax_cell_file):
        pic = slide.shapes.add_picture(minmax_cell_file, column1_x_1, row1_y_1, width=width_1, height=height_1)

###    Check if the polcurve image files exist     ###

#     for m in ['U1','U2','L2','L1']:
#         pol_curve_file = '{}/{}_pol_curve.png'.format(image_folder_name, m)
#         if os.path.exists(pol_curve_file):
#             pic = slide.shapes.add_picture(pol_curve_file, column2_x, row1_y, width=width, height=height)
    
    #if len(running_cell_positions) == 1:
        #del prs.slides._sldIdLst[6]
        
    elif len(running_cell_positions) == 2:
        keys_list = list(running_cell_positions.keys())  # Convert keys to a list
        slide = prs.slides[6]
        height = Inches(4.88)
        width = Inches(6.5)
        row1_y = Inches(1.39)
        column1_x = Inches(0.08)
        column2_x = Inches(6.76)
        
        pol_curve_file1 = '{}/{}_pol_curve.png'.format(image_folder_name, keys_list[0])
        pic = slide.shapes.add_picture(pol_curve_file1, column1_x, row1_y, width=width, height=height)
        pol_curve_file2 = '{}/{}_pol_curve.png'.format(image_folder_name, keys_list[1])
        pic = slide.shapes.add_picture(pol_curve_file2, column2_x, row1_y, width=width, height=height)

    elif len(running_cell_positions) == 3:
        keys_list = list(running_cell_positions.keys())
        slide = prs.slides[6]
        height = Inches(3.45)
        width = Inches(4.94)
        row1_y = Inches(0.48)
        row2_y = Inches(3.94)
        column1_x = Inches(0.13)
        column2_x = Inches(3.82)
        column3_x = Inches(7.5)
        
        pol_curve_file1 = '{}/{}_pol_curve.png'.format(image_folder_name, keys_list[0])
        pic = slide.shapes.add_picture(pol_curve_file1, column1_x, row1_y, width=width, height=height)
        pol_curve_file2 = '{}/{}_pol_curve.png'.format(image_folder_name, keys_list[1])
        pic = slide.shapes.add_picture(pol_curve_file2, column3_x, row1_y, width=width, height=height)
        pol_curve_file3 = '{}/{}_pol_curve.png'.format(image_folder_name, keys_list[2])
        pic = slide.shapes.add_picture(pol_curve_file3, column2_x, row2_y, width=width, height=height)
        
    elif len(running_cell_positions) == 4:
        keys_list = list(running_cell_positions.keys())
        slide = prs.slides[6]
        height = Inches(3.45)
        width = Inches(4.94)
        row1_y = Inches(0.52)
        row2_y = Inches(3.86)
        column1_x = Inches(0.13)
        column2_x = Inches(7.02)
        
        pol_curve_file1 = '{}/{}_pol_curve.png'.format(image_folder_name, keys_list[0])
        pic = slide.shapes.add_picture(pol_curve_file1, column1_x, row1_y, width=width, height=height)
        pol_curve_file2 = '{}/{}_pol_curve.png'.format(image_folder_name, keys_list[1])
        pic = slide.shapes.add_picture(pol_curve_file2, column2_x, row1_y, width=width, height=height)
        pol_curve_file3 = '{}/{}_pol_curve.png'.format(image_folder_name, keys_list[2])
        pic = slide.shapes.add_picture(pol_curve_file3, column1_x, row2_y, width=width, height=height)
        pol_curve_file3 = '{}/{}_pol_curve.png'.format(image_folder_name, keys_list[3])
        pic = slide.shapes.add_picture(pol_curve_file3, column2_x, row2_y, width=width, height=height)

        
    slide = prs.slides[7]
    height = Inches(5.89)
    width = Inches(8)
    row1_y = Inches(1.61)
    column1_x = Inches(0)
    # place the LWC Pressure 704 & 705 and PMP State chart
    pic = slide.shapes.add_picture('{}/LWC_Pressure_and_PMP_State.png'.format(image_folder_name), column1_x, row1_y,width=width, height=height)
    pic.click_action.hyperlink.address = '{}/LWC_Pressure_and_PMP_State.html'.format(image_folder_name)
    
    width=Inches(5.87)
    height=Inches(4.19)
    row2_y = Inches(3.03)
    column2_x = Inches(7.46)
    # place the lWC conductivity chart
    pic = slide.shapes.add_picture('{}/LWC_Conductivity.png'.format(image_folder_name), column2_x, row2_y,width=width, height=height)
    pic.click_action.hyperlink.address = '{}/LWC_Conductivity.html'.format(image_folder_name)

    # save the PowerPoint file with the customized name
    prs.save('{}/{}.pptx'.format(output_dir, ppt_fn))
    print("Presentation is published")


# In[29]:


def main(csv_file_list, stack_info_fn, output_dir, Unit_name):
    image_folder_name = '{}/images'.format(output_dir)
    if not os.path.exists(image_folder_name):
        os.mkdir(image_folder_name)
    
    if Unit_name == 'U05':
        template_ppt_fn = 'ckb_template_U05'
    elif Unit_name == 'U0002':
        template_ppt_fn = 'ckb_template_U0002'
    elif Unit_name == 'SPR001':
        template_ppt_fn = 'ntpc_template'
    ppt_fn = datetime.now().strftime("%Y-%B-%d_%H%M"+"hrs_IST_") + Unit_name +"_update"
    print("Output Rreport name will be:\n " + ppt_fn)
    ## giving user input for Stack position and Stackid    
    Stack_information = pd.read_excel(stack_info_fn, sheet_name = Unit_name, header=0)
#     print(csv_file_list)
    sys_data = pd.concat([pd.read_csv(f, on_bad_lines='warn') for f in csv_file_list ], ignore_index=True).fillna(np.nan)
    
    sys_data.to_csv('outputsample1.csv')

    #Calculating run hrs by considering timegap more than 10 sec as an interruption.
    # converting datetime from string to datetime format
    sys_data['Time'] = pd.to_datetime(sys_data['Time'],dayfirst=True)
    ## initialising run_hrs and a duration list
    run_hrs = 0
    duration_hr = [0.000]
    
    ## definig the run hrs completed by defining delta(s)    
    for i in range (1,len(sys_data)):
        s = (sys_data['Time'][i] - sys_data['Time'][i-1]).seconds
        if s <= 3600 :
            run_hrs = run_hrs + (s/3600)
            duration_hr.append(run_hrs)
        elif s > 3600:
            run_hrs = run_hrs + (10/3600)
            duration_hr.append(run_hrs)
    #adding a column for duration in dataframe
    sys_data['duration_hr'] = duration_hr        
    # sys_data[['Time','duration_hr']].head()
    # Output the total run duration in round figure in hrs (excluding interruptions)

    total_duration_hrs = round(max(sys_data['duration_hr']),0)
    print("Total Run Duration is "+ '"'+str(int(total_duration_hrs))+" hours"+'" (excluding interruptions)')
    Start_time = str(min(sys_data['Time']))
    End_time = str(max(sys_data['Time']))

    print("From : "+Start_time)
    print("To   : "+End_time)

    # replace the the string '-NA-' in data with blank
    sys_data = sys_data.replace('-NA-',np.nan)

    #remove rows that have any blank values
    # sys_data = sys_data.dropna(axis = 0)
    
    ## pre-processing the data ##
    ##remove the blank cols from data
    sys_data_empty = []
    for i in sys_data.columns:
        if len(sys_data[i].value_counts()) == 0:
            sys_data_empty.append(i)
            sys_data = sys_data.drop(i,axis = 1)        
    # print("The following fields are completely blank in the data, hence are removed & stored separately:", sys_data_empty)

    #convert the required fields from object to numeric datatype
    obj_cols = sys_data.loc[:, ~sys_data.columns.isin(['Time', 'MP State', 'MS State','Train 1 State', 'Train 2 State', 'PMP 102'])].columns
    sys_data[obj_cols] = sys_data[obj_cols].apply(pd.to_numeric, errors='coerce', axis=1)
    #make all the col names in uppercase
    sys_data.columns = sys_data.columns.str.upper()
    #rename TIME col as Date
    sys_data.rename(columns= {'TIME':'Date'},inplace = True)
    
    # Assign a unique order value based on the DataFrame's index
    sys_data['order'] = range(len(sys_data))
    
    # Sort the DataFrame by the 'order' column
    sys_data.sort_values(by=['order'], inplace=True)
    
    # dropping the cell v cols which are not active in the 112 stack design (every 9th and last 3) 
    # sys_data.drop(columns=['C8-14', 'C8-15', 'C8-16', 'C1-9','C2-2','C2-11','C3-4','C3-13','C4-6','C4-15','C5-8','C6-1','C6-10','C7-3','C7-12','C8-5'],inplace=True)
    # sys_data.drop(columns=['C16-14', 'C16-15', 'C16-16', 'C9-9', 'C10-2', 'C10-11', 'C11-4', 'C11-13', 'C12-6','C12-15', 'C13-8','C14-1','C14-10','C15-12','C15-12','C16-5'],inplace=True)

    # sys_data['Avg_V_C1_1_C8_16'] = (sys_data[List_C1_1_C8_16].sum(axis=1))/len(List_C1_1_C8_16)
    # sys_data['Max_V_C1_1_C8_16'] = sys_data[List_C1_1_C8_16].max(axis=1) 
    # sys_data['Min_V_C1_1_C8_16'] = sys_data[List_C1_1_C8_16].min(axis=1)
    
    # sys_data['Avg_V_C1_1-C8_16'] = (sys_data[List_C1_1_C8_16].sum(axis=1))/len(List_C1_1_C8_16)
    # sys_data['Max_V_C1_1-C8_16'] = sys_data[List_C1_1_C8_16].max(axis=1) 
    # sys_data['Min_V_C1_1-C8_16'] = sys_data[List_C1_1_C8_16].min(axis=1)    
    
    sys_data['CURRENT DENSITY'] = sys_data['STACK CURRENT'] / 702
    
    columns_to_check = sys_data.loc[:, 'C1-1':'C16-16'].columns


    # Assuming sys_data is your DataFrame
    columns_to_check = sys_data.loc[:, 'C1-1':'C16-16'].columns
    #sys_data['std_dev'] = sys_data[columns_to_check].std(axis=1)
        
# running the functions for creating dfs & charts
    Vcell_data = create_Vcell_data(sys_data)  
    plot_cd(sys_data, image_folder_name)
    #avg_min_max_stack_v(sys_data, image_folder_name)
    #avg_min_max_stack_v_old_singleaxis(sys_data, image_folder_name)
    avg_min_max_stack_v_Ben(Vcell_data, sys_data, image_folder_name)
    #avg_min_max_stack_v_old(sys_data, image_folder_name)
    # ASR_chart(Vcell_data_asr, image_folder_name)
    sys_data = limit_setting(sys_data)
    pressure_and_conductivity(sys_data, image_folder_name)
    ttc_ktc_lvl_sensors(sys_data, image_folder_name)
    #HYS102_vs_LVL101(sys_data, image_folder_name)
    hyd_and_water_sensor(sys_data, image_folder_name)
    LWC_Pressure_and_PMP_State(sys_data, image_folder_name)
    LWC_conductivity(sys_data, image_folder_name)
    T1, T2, N = t_start_t_end_N(Vcell_data)
    stack_specific_energy(Vcell_data, image_folder_name)
    #stack_specific_energy_new(sys_data, image_folder_name)
    
    #if T1/T2 is empty, then do not run pol curve related functions else run them
    if len(T1) == 0:
        pass
        print('No pol curve in the data')
        running_cell_positions = {}  # since it is an input in publish ppt function, it has to be defined
    else:
        cd_and_polcurve(Vcell_data, T1, T2, image_folder_name)
        N_temp, T1_temp, T2_temp = avg_stack_v_polcurve(N, T1, T2, Vcell_data, image_folder_name)
        if len(T1_temp) == 0:
            print('Pol curve not found in the data') 
        else:
            MinMax_v_cell(N_temp, T1_temp, T2_temp, Vcell_data, image_folder_name)
            running_cell_positions = stack_wise_polcurve(N, T1, T2, Vcell_data, image_folder_name)
    
#   running the function to publish ppt
    publish_presentation(template_ppt_fn, image_folder_name, ppt_fn, output_dir, total_duration_hrs, Start_time, End_time, Stack_information, running_cell_positions)  

    
# In[30]:


def rungui():
    
    # Create an instance of tkinter frame
    win = Tk()
    win.title('Lotus Sys Data Visualization')
    
    # create the frame layout for the main window
    frame = tk.Frame(master=win, width=500, height=175)
    frame.pack()
    
    # add a label for the automatic data
    label1 = tk.Label(master=frame, text="Data files:")
    label1.place(x=65, y=25)
    
    # add an entry box for the automatic data file
    entry1 = tk.Entry(win, bg="white", width = 45)
    entry1.place(x=125, y=25)
    
    # this function opens a file browser for excel files and outputs the selected box to the automatic data entry field
    def browse1():
        try:
            entry1.delete(0, tk.END)
            files = filedialog.askopenfilenames(filetypes=[('CSV Files', '*.csv')])
            for i in files[:-1]:
                entry1.insert(tk.END, '{}, '.format(i))
            entry1.insert(tk.END, files[-1])
        except PermissionError:
            messagebox.showerror('Permission Error', 'Error: Please close the file on your computer and try again.')

    # Create a browse button that links to the browse1 function
    button1 = ttk.Button(win, text="Browse", command=browse1)
    button1.place(x=400, y=21)

    # create a label for the manual data
    label2 = tk.Label(master=frame, text="Stack Information file:")
    label2.place(x=0, y=50)

    # create an entry box for the manual file name
    entry2 = tk.Entry(win, bg="white", width = 45)
    entry2.place(x=125, y=50)

    # this function opens a file browser for excel files and outputs the 
    # selected box to the manual data entry field
    def browse2():
        try:
            file = filedialog.askopenfile(mode='r', filetypes=[('Excel Files', 
                                                                '*.xlsx')])
            filename = os.path.abspath(file.name)
            entry2.insert(tk.END, filename)
        except PermissionError:
            messagebox.showerror('Permission Error', 'Error: Please close the file on your computer and try again.')

    # Create a browse button that links to the browse2 function
    button2 = ttk.Button(win, text="Browse", command=browse2)
    button2.place(x=400, y=47)

    #create a label for the output folder
    label3 = tk.Label(master=frame, text="Output folder:")
    label3.place(x=40, y=77)

    # creat an entry box for the output folder    
    entry3 = tk.Entry(win, bg="white", width = 45)
    entry3.place(x=125, y=77)

    # this function opens a file browser for excel files and outputs the 
    # selected box to the manual data entry field
    def browse3():
        try:
            folder = filedialog.askdirectory()
            #foldername = os.path.abspath(folder.name)
            entry3.insert(tk.END, folder)
        except PermissionError:
            messagebox.showerror('Permission Error', 'Error: Please close the folder on your computer and try again.')

    # Create a browse button that links to the browse3 function
    button3 = ttk.Button(win, text="Browse", command=browse3)
    button3.place(x=400, y=74)
    
    #create a radio item for selecting the template
    label4 = tk.Label(master=frame, text="System name:")
    label4.place(x=0, y=100)
    selected_system = tk.StringVar()
    radio1 = Radiobutton(win, text='U05', value='U05', variable=selected_system)
    radio1.place(x=85, y=100)
    radio2 = Radiobutton(win, text='U0002', value='U0002', variable=selected_system)
    radio2.place(x=140, y=100)
    radio3 = Radiobutton(win, text='SPR001', value='SPR001', variable=selected_system)
    radio3.place(x=210, y=100)
    
    # this defines the main run function that we wish to do
    def run_main():
        template_ppt_fn = selected_system.get()
        csv_file_list = entry1.get().replace('.csv,' , '.csv;').split('; ')
        stack_info_fn = entry2.get()
        output_dir = entry3.get()
        main(csv_file_list, stack_info_fn, output_dir, template_ppt_fn)
        win.destroy()
        

    # this is the overall function called that does the main function and the progress bar
    def run_function(name, func):
        # disable the run button
        button4['state'] = 'disabled'
        # define and start the progress bar
        progress_bar = ttk.Progressbar(win, orient = HORIZONTAL, length = 400, mode='indeterminate')
        progress_bar.place(x=50, y = 125)
        progress_bar.start(interval = 15)
        # run the main function
        func()
        # stop the progress bar
        progress_bar.stop()
        # print program complete
        print('Program complete')
        #destory the progress bar
        progress_bar.destroy()
        # show the finished label
        label4 = tk.Label(master=frame, text="Finished!")
        label4.place(x=240, y=115)
        # renable the run button
        button4['state'] = 'enabled'

    # this is the function that threads together the above two functions
    def run_thread(name, func):
        Thread(target = run_function, args=(name,func)).start()

    # this is the function that is called when the run button is clicked
    def run_clicked():
        run_thread('main', run_main)

    # Create a Button that runs the main script
    button4 = ttk.Button(win, text="Run", command=run_clicked)
    button4.place(x=225, y=145)

    # create a help menu

    win.mainloop()
#%%    
rungui()   

# csv_file_list = [  'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/extracted_files_ntpc/lhc1_2023_01_25.csv'

#                    'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/extracted_files_ntpc/lhc1_2024_01_02.csv',
#                    'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/extracted_files_ntpc/lhc1_2024_01_03.csv',
#                    'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/extracted_files_ntpc/lhc1_2024_01_08.csv',
#                    'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/extracted_files_ntpc/lhc1_2024_01_09.csv',
#                    'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/extracted_files_ntpc/lhc1_2024_01_10.csv',
#                    'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/extracted_files_ntpc/lhc1_2024_01_11.csv'
#     ]
# stack_info_fn = 'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/Stack_Information.xlsx'
# output_dir = 'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/updates'
# Unit_name = 'SPR001'

# # main(csv_file_list, stack_info_fn, output_dir, unit_name)
# #%%
# csv_file_list = [
#                   'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/extracted_files_ntpc/lhc1_2023_12_02.csv',
#                   'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/extracted_files_ntpc/lhc1_2023_12_03.csv',
#                   'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/extracted_files_ntpc/lhc1_2023_12_04.csv',
#                   'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/extracted_files_ntpc/lhc1_2023_12_05.csv',
#                   'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/extracted_files_ntpc/lhc1_2023_12_06.csv',
#                   'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/extracted_files_ntpc/lhc1_2023_12_07.csv',
#                   'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/extracted_files_ntpc/lhc1_2023_12_08.csv',
#                   # 'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/extracted_files_ntpc/lhc1_2023_12_09.csv',
#                   # 'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/extracted_files_ntpc/lhc1_2023_12_10.csv',
#                   'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/extracted_files_ntpc/lhc1_2023_12_11.csv',
#                   'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/extracted_files_ntpc/lhc1_2023_12_12.csv',
#                   'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/extracted_files_ntpc/lhc1_2023_12_13.csv',
#                   'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/extracted_files_ntpc/lhc1_2023_12_14.csv',
#                   'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/extracted_files_ntpc/lhc1_2023_12_15.csv',
#                   'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/extracted_files_ntpc/lhc1_2023_12_16.csv',
#                   'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/extracted_files_ntpc/lhc1_2023_12_17.csv',
#                   'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/extracted_files_ntpc/lhc1_2023_12_18.csv',
#                   'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/extracted_files_ntpc/lhc1_2023_12_19.csv',
#                   'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/extracted_files_ntpc/lhc1_2023_12_20.csv',
#                   'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/extracted_files_ntpc/lhc1_2023_12_21.csv',
#                   'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/extracted_files_ntpc/lhc1_2023_12_22.csv',
#                   'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/extracted_files_ntpc/lhc1_2023_12_23.csv',
#                   'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/extracted_files_ntpc/lhc1_2023_12_24.csv',
#                   'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/extracted_files_ntpc/lhc1_2023_12_25.csv',
#                   'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/extracted_files_ntpc/lhc1_2023_12_26.csv',
#                   'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/extracted_files_ntpc/lhc1_2023_12_27.csv',
#                   'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/extracted_files_ntpc/lhc1_2023_12_28.csv',


#                   'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/extracted_files_ntpc/lhc1_2024_01_01.csv',
#                   'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/extracted_files_ntpc/lhc1_2024_01_02.csv',
#                   'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/extracted_files_ntpc/lhc1_2024_01_03.csv',
#                   'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/extracted_files_ntpc/lhc1_2024_01_04.csv',
#                   'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/extracted_files_ntpc/lhc1_2024_01_08.csv',
#                   'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/extracted_files_ntpc/lhc1_2024_01_09.csv',
#                   'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/extracted_files_ntpc/lhc1_2024_01_10.csv',
# 'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/extracted_files_ntpc/lhc1_2024_01_11.csv'
# ]
# stack_info_fn = 'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/Stack_Information.xlsx'
# output_dir = 'C:/Users/nidhi.dubey/OneDrive - Ohmium International, Inc/Documents/Lotus_system_data/updates'
# Unit_name = 'SPR001'

# #%%
# image_folder_name = '{}/images'.format(output_dir)
# if not os.path.exists(image_folder_name):
#     os.mkdir(image_folder_name)

# if Unit_name == 'U05':
#     template_ppt_fn = 'ckb_template_U05'
# elif Unit_name == 'U0002':
#     template_ppt_fn = 'ckb_template_U0002'
# elif Unit_name == 'SPR001':
#     template_ppt_fn = 'ntpc_template'
# ppt_fn = datetime.now().strftime("%Y-%B-%d_%H%M"+"hrs_IST_") + Unit_name +"_update"
# print("Output Rreport name will be:\n " + ppt_fn)
# ## giving user input for Stack position and Stackid    
# Stack_information = pd.read_excel(stack_info_fn, sheet_name = Unit_name, header=0)
# #     print(csv_file_list)
# sys_data = pd.concat([pd.read_csv(f, on_bad_lines='warn') for f in csv_file_list ], ignore_index=True).fillna(np.nan)

# sys_data.to_csv('outputsample1.csv')

# #Calculating run hrs by considering timegap more than 10 sec as an interruption.
# # converting datetime from string to datetime format
# sys_data['Time'] = pd.to_datetime(sys_data['Time'],dayfirst=True)
# ## initialising run_hrs and a duration list
# run_hrs = 0
# duration_hr = [0.000]

# ## definig the run hrs completed by defining delta(s)    
# for i in range (1,len(sys_data)):
#     s = (sys_data['Time'][i] - sys_data['Time'][i-1]).seconds
#     if s <= 3600 :
#         run_hrs = run_hrs + (s/3600)
#         duration_hr.append(run_hrs)
#     elif s > 3600:
#         run_hrs = run_hrs + (10/3600)
#         duration_hr.append(run_hrs)
# #adding a column for duration in dataframe
# sys_data['duration_hr'] = duration_hr        
# # sys_data[['Time','duration_hr']].head()
# # Output the total run duration in round figure in hrs (excluding interruptions)

# total_duration_hrs = round(max(sys_data['duration_hr']),0)
# print("Total Run Duration is "+ '"'+str(int(total_duration_hrs))+" hours"+'" (excluding interruptions)')
# Start_time = str(min(sys_data['Time']))
# End_time = str(max(sys_data['Time']))

# print("From : "+Start_time)
# print("To   : "+End_time)

# # replace the the string '-NA-' in data with blank
# sys_data = sys_data.replace('-NA-',np.nan)

# #remove rows that have any blank values
# # sys_data = sys_data.dropna(axis = 0)

# ## pre-processing the data ##
# ##remove the blank cols from data
# sys_data_empty = []
# for i in sys_data.columns:
#     if len(sys_data[i].value_counts()) == 0:
#         sys_data_empty.append(i)
#         sys_data = sys_data.drop(i,axis = 1)        
# # print("The following fields are completely blank in the data, hence are removed & stored separately:", sys_data_empty)

# #convert the required fields from object to numeric datatype
# obj_cols = sys_data.loc[:, ~sys_data.columns.isin(['Time', 'MP State', 'MS State','Train 1 State', 'Train 2 State', 'PMP 102'])].columns
# sys_data[obj_cols] = sys_data[obj_cols].apply(pd.to_numeric, errors='coerce', axis=1)
# #make all the col names in uppercase
# sys_data.columns = sys_data.columns.str.upper()
# #rename TIME col as Date
# sys_data.rename(columns= {'TIME':'Date'},inplace = True)

# # Assign a unique order value based on the DataFrame's index
# sys_data['order'] = range(len(sys_data))

# #create a desired new column with only day & month
# sys_data['date_&_hrs'] = sys_data['Date'].dt.strftime('%d-%b'+'<br>'+'%H'+':00')

# # Convert the 'date_&_hrs' column to datetime
# #sys_data['date_&_hrs'] = pd.to_datetime(sys_data['date_&_hrs'], format='%d-%b %H:%M')
# sys_data['date_&_hrs'] = pd.to_datetime(sys_data['date_&_hrs'], format='%d-%b'+'<br>'+'%H'+':00')

# # Sort the DataFrame by the 'date_&_hrs' column and then by the 'order' column
# #sys_data.sort_values(by=['date_&_hrs', 'order'], inplace=True)
# sys_data.sort_values(by=['order'], inplace=True)


# # dropping the cell v cols which are not active in the 112 stack design (every 9th and last 3) 
# sys_data.drop(columns=['C8-14', 'C8-15', 'C8-16', 'C1-9','C2-2','C2-11','C3-4','C3-13','C4-6','C4-15','C5-8','C6-1','C6-10','C7-3','C7-12','C8-5'],inplace=True)

# sys_data['Avg_V_C1_1-C8_16'] = (sys_data[List_C1_1_C8_16].sum(axis=1))/len(List_C1_1_C8_16)
# sys_data['Max_V_C1_1-C8_16'] = sys_data[List_C1_1_C8_16].max(axis=1) 
# sys_data['Min_V_C1_1-C8_16'] = sys_data[List_C1_1_C8_16].min(axis=1)
# sys_data['CURRENT DENSITY'] = sys_data['STACK CURRENT'] / 702
# sys_data['std_dev'] = sys_data.loc[:,Stack_voltage_list].std(axis=1)

# #sys_data.to_csv('outputsample2.csv')

# # #%%
# #def avg_V_std_dev_errorbands(sys_data,image_folder_name):
# mean_trace = go.Scatter(
#         x = sys_data.Date,
#         y = sys_data['Avg_V_C1_1-C8_16'],
#         mode = 'lines',
#         name = 'Avg V',
#         #line = dict(color = 'rgb(31, 119, 180)', shape='hv'),
#         line=dict(color='rgb(31, 119, 180)', shape='hv', width=0.5),
#         connectgaps=False
#         )
# max_trace = go.Scatter(
#     x=list(sys_data.Date), 
#     y=list(sys_data['Max_V_C1_1-C8_16']), 
#     line=dict(width=0, shape = 'hv', color = 'red'),
#     mode='lines',
#     name = 'max',
#     showlegend=False,
#     connectgaps=False
#     )
# min_trace = go.Scatter(
#     name='Min',
#     x=list(sys_data.Date),
#     y=list(sys_data['Min_V_C1_1-C8_16']),
#     mode='lines',
#     fillcolor='rgba(31,119,180,0.2)',
#     line=dict(width = 0, shape = 'hv', color = 'green'),
#     fill='tonexty',
#     showlegend=False,
#     connectgaps=False
# )
 
# std_dev_trace = go.Scatter(
#     x = sys_data.Date,
#     y = sys_data['std_dev'],
#     mode = 'lines',
#     name = 'Std Dev',
#     #line = dict(color = 'rgb(255, 127, 14)', shape='hv'),
#     line=dict(color='rgb(255, 127, 14)', shape='hv', width=0.5),
#     connectgaps=False
#     )
 
# dashed_line = {
#         'type': 'line',
#         'xref': 'paper',
#         'yref': 'y',
#         'x0': 0,
#         'y0': 2,
#         'x1': 1,
#         'y1': 2,
#         'line': {
#             'color': 'black',
#             'width': 2,
#             'dash': 'dash',
#         },
#     }
# data = [mean_trace, max_trace, min_trace, std_dev_trace]
 
# layout = go.Layout(
#             #title = 'plot_title',
#             #to define axis, can either use go.layout.etc or dict(etc)
#             # xaxis = go.layout.XAxis(title = 'Z<sub>R</sub> [mOhm.cm<sup>2</sup>]',
#             #                        showgrid = False),
#             #yaxis = dict(title='',
#             #             showgrid = False),
#             # legend= dict(x=0, y =1.1, bgcolor = 'rgba(0,0,0,0)'),
#             # legend_title = None,\
#             tickformat = "%d-%b",
#             hoverformat = "%d-%b %H:00",    
#             template = 'simple_white',
#             shapes=[dashed_line]
#             # font = dict(size = 18)
#             )
    
# fig = go.Figure(data = data, layout = layout)
# fig.show()
# # fn = "C:/Users/Ben/OneDrive - Ohmium International/Documents/Sandbox/Spirare_Report_Draft/Charts/{}.html".format(plot_title.replace(' ', '_'))
# # fig.write_html(fn)

# fig.write_image('{}/LHC_Stack_Voltage_&_Current_Density.png'.format(image_folder_name), engine="orca")
# fig.write_html('{}/LHC_Stack_Voltage_&_Current_Density.html'.format(image_folder_name), include_plotlyjs=('cdn'))
